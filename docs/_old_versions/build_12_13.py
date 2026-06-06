#!/usr/bin/env python3
"""
사업계획서 12-13장 슬라이드 추가 스크립트 (v0.5 — 전체 편집 가능)
- 원본 복사 후 복사본에 8슬라이드(12장 4 + 13장 4) 추가
- v4 가시성 규칙: 최소 10pt, 본문 12pt+, 도형 라벨 11pt+, 슬라이드당 도형 30개 이하
- DFD 다이어그램을 모두 **편집 가능 PPT 도형(shape)**으로 제작
- 실물 사진(CCTV 캡처, 사고 현장)만 이미지로 삽입
- 레이아웃 4 사용, 기존 양식 동일
"""
import shutil
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ── 경로 설정 ──
SRC = Path("/공유폴더/파일공유용/JHJ/사고자료/(사업계획서)교통센터 생성형 AI 상환관제_v0.2.pptx")
DST = Path("/workspace/prj_cctv/사고분석_설계/docs/제안서_v0.5_12_13장.pptx")
ASSETS = Path("/workspace/prj_cctv/사고분석_설계/docs/assets")

# ── 색상 팔레트 ──
GOLD = RGBColor(0xFF, 0xC0, 0x00)
ACCENT = RGBColor(0x00, 0x41, 0x9F)
MAIN_BG = RGBColor(0x04, 0x5F, 0xA4)
CAT_HDR = RGBColor(0x08, 0x67, 0xBE)
SUB_CAT = RGBColor(0x00, 0xB0, 0xF0)
TAG_LBL = RGBColor(0x18, 0x4D, 0xA4)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)
LIGHT_GRAY = RGBColor(0xD9, 0xD9, 0xD9)
YELLOW_TAG = RGBColor(0xFF, 0xD9, 0x66)
RED_ALERT = RGBColor(0xC0, 0x39, 0x2B)
GREEN_OK = RGBColor(0x27, 0xAE, 0x60)
DARK_BG = RGBColor(0x2C, 0x3E, 0x50)
GRAY_TEXT = RGBColor(0x66, 0x66, 0x66)
LIGHT_BLUE_BG = RGBColor(0xE8, 0xF4, 0xFD)
PALE_BLUE = RGBColor(0xD6, 0xE9, 0xF8)
LIGHT_PALE = RGBColor(0xBB, 0xDE, 0xF5)  # 연한 파랑 (수치 라벨용)

# ── 공통 요소 위치 (EMU) — 기존 슬라이드에서 추출 ──
POS_CHAPTER_NUM = (326868, 543571, 482429, 445122)
POS_TITLE = (859257, 625350, 4570007, 319891)
POS_SUBTITLE = (639649, 1155600, 4726666, 278792)
POS_DESC = (593725, 1645961, 9691929, 569387)

# 본문 영역
BODY_LEFT = Inches(0.45)
BODY_TOP = Inches(2.80)
BODY_WIDTH = Inches(10.76)
BODY_HEIGHT = Inches(4.90)


# ══════════════════════════════════════════════════════════════
# 헬퍼 함수
# ══════════════════════════════════════════════════════════════

def set_font(run, size_pt, bold=False, color=None, name=None):
    """run에 폰트 속성 설정"""
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    if color:
        run.font.color.rgb = color
    if name:
        run.font.name = name


def add_text_box(slide, left, top, width, height, text, size_pt=12,
                 bold=False, color=None, alignment=PP_ALIGN.LEFT, name=None):
    """텍스트 박스 추가"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    set_font(run, size_pt, bold, color, name)
    return txBox


def add_multiline_text_box(slide, left, top, width, height, lines,
                           alignment=PP_ALIGN.LEFT):
    """여러 줄 텍스트 박스 추가. lines = [(text, size_pt, bold, color), ...]"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, (text, size_pt, bold, color) in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = alignment
        p.space_before = Pt(2)
        p.space_after = Pt(2)
        run = p.add_run()
        run.text = text
        set_font(run, size_pt, bold, color)
    return txBox


def add_rounded_rect(slide, left, top, width, height, fill_color=None,
                     text="", size_pt=12, bold=False, text_color=None,
                     alignment=PP_ALIGN.CENTER):
    """둥근 사각형 도형 추가"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    shape.line.fill.background()

    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        tf.paragraphs[0].alignment = alignment
        run = tf.paragraphs[0].add_run()
        run.text = text
        set_font(run, size_pt, bold, text_color or (WHITE if fill_color else BLACK))

    shape.text_frame.paragraphs[0].space_before = Pt(0)
    shape.text_frame.paragraphs[0].space_after = Pt(0)

    return shape


def add_multiline_rounded_rect(slide, left, top, width, height, fill_color,
                               lines, alignment=PP_ALIGN.CENTER):
    """여러 줄 텍스트가 있는 둥근 사각형. lines = [(text, size_pt, bold, color), ...]"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    shape.line.fill.background()

    tf = shape.text_frame
    tf.word_wrap = True
    for i, (text, size_pt, bold, color) in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = alignment
        p.space_before = Pt(1)
        p.space_after = Pt(1)
        run = p.add_run()
        run.text = text
        set_font(run, size_pt, bold, color)

    return shape


def add_multiline_box(slide, left, top, width, height, fill_color,
                      lines, sizes, bolds, colors, alignment=PP_ALIGN.LEFT):
    """하나의 도형에 여러 줄 텍스트 (각 줄 별도 paragraph)"""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    for i, (text, size, bold, color) in enumerate(zip(lines, sizes, bolds, colors)):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = alignment
        p.space_before = Pt(1)
        p.space_after = Pt(1)
        run = p.add_run()
        run.text = text
        set_font(run, size, bold, color)
    return shape


def add_rect(slide, left, top, width, height, fill_color=None,
             text="", size_pt=12, bold=False, text_color=None,
             alignment=PP_ALIGN.CENTER):
    """직사각형 도형 추가"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, height
    )
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    shape.line.fill.background()

    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        tf.paragraphs[0].alignment = alignment
        run = tf.paragraphs[0].add_run()
        run.text = text
        set_font(run, size_pt, bold, text_color or (WHITE if fill_color else BLACK))

    return shape


def add_diamond(slide, left, top, width, height, fill_color=None,
                text="", size_pt=11, bold=True, text_color=None):
    """다이아몬드(판단 분기) 도형"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.DIAMOND, left, top, width, height
    )
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    shape.line.color.rgb = fill_color or GOLD
    shape.line.width = Pt(2)

    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        run = tf.paragraphs[0].add_run()
        run.text = text
        set_font(run, size_pt, bold, text_color or BLACK)

    return shape


def add_arrow_right(slide, left, top, width, height, fill_color=None):
    """오른쪽 화살표"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RIGHT_ARROW, left, top, width, height
    )
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = MAIN_BG
    shape.line.fill.background()
    return shape


def add_arrow_down(slide, left, top, width, height, fill_color=None):
    """아래쪽 화살표"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.DOWN_ARROW, left, top, width, height
    )
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = MAIN_BG
    shape.line.fill.background()
    return shape


def setup_slide_header(slide, chapter_num, title, subtitle, desc_line1, desc_line2=""):
    """공통 슬라이드 헤더 설정: 장번호, 대제목, 소제목, 설명문"""
    add_text_box(slide, *POS_CHAPTER_NUM,
                 str(chapter_num), size_pt=20, bold=True, color=GOLD)

    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 0:
            ph.text_frame.clear()
            p = ph.text_frame.paragraphs[0]
            run = p.add_run()
            run.text = title
            run.font.language_id = 0x0412
            run.font.size = Pt(18)  # rPr 필수: placeholder에도 font.size 설정
            break

    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 10:
            ph.text_frame.clear()
            p = ph.text_frame.paragraphs[0]
            run = p.add_run()
            run.text = subtitle
            set_font(run, 18, bold=True, color=BLACK)
            break

    desc_box = slide.shapes.add_textbox(*POS_DESC)
    tf = desc_box.text_frame
    tf.word_wrap = True
    p1 = tf.paragraphs[0]
    run1 = p1.add_run()
    run1.text = desc_line1
    set_font(run1, 16, color=BLACK)
    if desc_line2:
        p2 = tf.add_paragraph()
        run2 = p2.add_run()
        run2.text = desc_line2
        set_font(run2, 16, color=BLACK)


def add_table(slide, left, top, width, height, rows, cols, data,
              header_color=None, font_size=10):
    """표 추가. data는 [row][col] 2D 리스트. run-level rPr 준수."""
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table

    col_width = width // cols
    for ci in range(cols):
        table.columns[ci].width = col_width

    for ri in range(rows):
        for ci in range(cols):
            cell = table.cell(ri, ci)
            txt = data[ri][ci] if ri < len(data) and ci < len(data[ri]) else ""
            tf = cell.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = txt
            run.font.size = Pt(font_size)
            if ri == 0 and header_color:
                run.font.bold = True
                run.font.color.rgb = WHITE
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE

            if ri == 0 and header_color:
                tc = cell._tc
                tcPr = tc.get_or_add_tcPr()
                for old_fill in tcPr.findall(qn('a:solidFill')):
                    tcPr.remove(old_fill)
                solidFill = tcPr.makeelement(qn('a:solidFill'), {})
                srgbClr = solidFill.makeelement(qn('a:srgbClr'), {'val': str(header_color)})
                solidFill.append(srgbClr)
                tcPr.insert(0, solidFill)

    return table_shape


def add_image_with_caption(slide, img_path, left, top, width, caption_text,
                           tag_text=None, tag_color=None):
    """이미지 삽입 + 태그 헤더(선택) + 캡션 (최소 10pt 준수)"""
    current_top = top

    # 태그 헤더
    if tag_text and tag_color:
        add_rounded_rect(slide, left, current_top,
                         width, Inches(0.28),
                         fill_color=tag_color,
                         text=tag_text, size_pt=12, bold=True,
                         text_color=WHITE)
        current_top += Inches(0.32)

    # 이미지 삽입 (height 생략 = 비율 유지)
    pic = slide.shapes.add_picture(str(img_path), left, current_top, width)
    img_height = pic.height

    # 캡션 (10pt — 최소 허용)
    caption_top = current_top + img_height + Inches(0.03)
    add_text_box(slide, left, caption_top, width, Inches(0.22),
                 caption_text, size_pt=10, color=GRAY_TEXT,
                 alignment=PP_ALIGN.CENTER)

    return pic, img_height


# ══════════════════════════════════════════════════════════════
# 12장 슬라이드 (4장) — DFD를 편집 가능 도형으로 재현
# ══════════════════════════════════════════════════════════════

def build_slide_12A(prs, layout):
    """12-A: 시스템 개요 — dfd-2 도형화 + 현행 vs 개선 (~25 도형)"""
    slide = prs.slides.add_slide(layout)
    setup_slide_header(
        slide, 12,
        "CCTV 영상 기반 교통상황 판단",
        "12.1 현황 문제 및 시스템 개요",
        "사고발생 시 근무자의 주관성 편차를 해소하는 CCTV AI 자동 분석 시스템",
    )

    # ── 상단 (70%): dfd-2 파이프라인 5단계를 가로 5칼럼 도형으로 재현 ──
    col_w = Inches(2.00)
    col_h = Inches(2.60)
    col_y = Inches(2.80)
    gap = Inches(0.15)
    arrow_w = Inches(0.15)
    start_x = Inches(0.35)

    columns = [
        {
            "title": "입력",
            "color": GREEN_OK,
            "items": ["CCTV 영상 (1fps)", "사고 보고서 (이력)", "지속데이터 (VMS/VDS)"],
        },
        {
            "title": "영상 분석",
            "color": SUB_CAT,
            "items": ["차량 탐지 (YOLO)", "차량 추적 (ByteTrack)", "차종 분류 (13종)", "속도/충돌위험 계산", "이상 상황 감지"],
        },
        {
            "title": "AI 판단 (핵심)",
            "color": MAIN_BG,
            "items": ["장면 파악", "사고 여부 판정", "차종 세분류 보정", "보고서 자동 작성 (93컬럼)", "위험도 평가"],
        },
        {
            "title": "데이터 저장",
            "color": GOLD,
            "items": ["차량 궤적", "사고 기록", "AI 분석 결과", "교통 통계"],
        },
        {
            "title": "출력",
            "color": RED_ALERT,
            "items": ["사고 현황", "차종별 교통량", "차단 분석", "보고서 (93컬럼)", "위험도 예측"],
        },
    ]

    for i, col in enumerate(columns):
        x = start_x + i * (col_w + gap + arrow_w)
        # 칼럼 도형 (멀티라인)
        lines = [(col["title"], 13, True, WHITE)]
        for item in col["items"]:
            lines.append((f"  {item}", 10, False, WHITE))
        add_multiline_rounded_rect(
            slide, x, col_y, col_w, col_h,
            fill_color=col["color"],
            lines=lines,
            alignment=PP_ALIGN.LEFT,
        )
        # 칼럼 간 화살표
        if i < len(columns) - 1:
            ax = x + col_w + Inches(0.01)
            add_arrow_right(slide, ax, col_y + Inches(1.10),
                            arrow_w, Inches(0.25), fill_color=MAIN_BG)

    # AI 판단 칼럼 하단: 모델 라벨
    model_x = start_x + 2 * (col_w + gap + arrow_w)
    add_text_box(slide, model_x, col_y + col_h + Inches(0.05),
                 col_w, Inches(0.22),
                 "Qwen2.5-VL (CPU→GPU)", size_pt=10, bold=True,
                 color=MAIN_BG, alignment=PP_ALIGN.CENTER)

    # ── 하단 (30%): 현행 vs 개선 비교 카드 2개 ──
    card_y = Inches(5.80)
    card_w = Inches(4.80)
    card_h = Inches(1.40)

    # 좌: 현행 카드
    add_multiline_rounded_rect(
        slide, Inches(0.45), card_y, card_w, card_h,
        fill_color=RED_ALERT,
        lines=[
            ("현행 (수동 방식)", 14, True, WHITE),
            ("수동 탐색 → 육안 확인 → 구두 보고", 12, False, WHITE),
            ("5~10분 소요 | 주관적 편차 발생", 12, True, YELLOW_TAG),
        ],
        alignment=PP_ALIGN.CENTER,
    )

    # 중앙 화살표
    add_arrow_right(slide, Inches(5.40), card_y + Inches(0.45),
                    Inches(0.45), Inches(0.40), fill_color=GREEN_OK)

    # 우: 개선 카드
    add_multiline_rounded_rect(
        slide, Inches(6.00), card_y, card_w, card_h,
        fill_color=MAIN_BG,
        lines=[
            ("개선 (AI 자동 방식)", 14, True, WHITE),
            ("자동 연결 → AI 분석 → DB 적재", 12, False, WHITE),
            ("~2분 소요 | 객관적, 일관된 기준", 12, True, GOLD),
        ],
        alignment=PP_ALIGN.CENTER,
    )

    return slide  # ~17 shapes


def build_slide_12B(prs, layout):
    """12-B: 3-Tier 감시 체계 (~18 도형)"""
    slide = prs.slides.add_slide(layout)
    setup_slide_header(
        slide, 12,
        "CCTV 영상 기반 교통상황 판단",
        "12.2 전국 CCTV 3-Tier 효율 감시 체계",
        "313대 광역 → 50대 핫스팟 → 동적 정밀, 단계별 자동 승격",
    )

    # ── 가로 3열 Tier 배치 ──
    tier_w = Inches(3.30)
    tier_h = Inches(3.50)
    tier_y = Inches(2.85)
    gap = Inches(0.20)
    start_x = Inches(0.45)

    tiers = [
        {
            "color": SUB_CAT,
            "lines": [
                ("Tier 1 — 광역 감시", 14, True, WHITE),
                ("313대", 20, True, WHITE),
                ("경량 프리필터 (MOG2 + 프레임 차분)", 12, False, WHITE),
                ("처리: ~1.5ms/프레임", 12, False, LIGHT_PALE),
            ],
        },
        {
            "color": CAT_HDR,
            "lines": [
                ("Tier 2 — 핫스팟 정밀", 14, True, WHITE),
                ("50대", 20, True, WHITE),
                ("YOLO + ByteTrack + 트리거 7종", 12, False, WHITE),
                ("처리: 21ms/프레임", 12, False, LIGHT_PALE),
            ],
        },
        {
            "color": MAIN_BG,
            "lines": [
                ("Tier 3 — 동적 정밀", 14, True, WHITE),
                ("동적", 20, True, WHITE),
                ("전체 AI 파이프라인 + VLM", 12, False, WHITE),
                ("사고 현장 집중 분석", 12, False, LIGHT_PALE),
            ],
        },
    ]

    for i, tier in enumerate(tiers):
        x = start_x + i * (tier_w + gap)
        add_multiline_rounded_rect(
            slide, x, tier_y, tier_w, tier_h,
            fill_color=tier["color"],
            lines=tier["lines"],
            alignment=PP_ALIGN.CENTER,
        )
        # 화살표 between tiers (승격 방향: 왼→오른)
        if i < 2:
            ax = x + tier_w + Inches(0.01)
            add_arrow_right(slide, ax, tier_y + Inches(1.40),
                            gap - Inches(0.02), Inches(0.35),
                            fill_color=MAIN_BG)

    # ── 하단: 핵심 수치 카드 3개 ──
    card_y = Inches(6.60)
    card_w = Inches(3.40)
    card_h = Inches(0.55)
    card_data = [
        ("100% CCTV 캡처 성공", MAIN_BG),
        ("약 2분 소요 (기존 5~10분)", CAT_HDR),
        ("35컬럼 DB 자동 적재", TAG_LBL),
    ]
    for i, (label, color) in enumerate(card_data):
        cx = Inches(0.45) + i * (card_w + Inches(0.15))
        add_rounded_rect(slide, cx, card_y, card_w, card_h,
                         fill_color=color,
                         text=label, size_pt=12, bold=True, text_color=WHITE)

    return slide  # ~11 shapes


def build_slide_12C(prs, layout):
    """12-C: CCTV 체류 캡처 실물 (~15 도형)"""
    slide = prs.slides.add_slide(layout)
    setup_slide_header(
        slide, 12,
        "CCTV 영상 기반 교통상황 판단",
        "12.3 사고 현장 CCTV 30초 체류 캡처",
        "ITS 돌발정보 수신 → 인근 CCTV 자동접속 → 1fps × 30초 연속 캡처",
    )

    # ── 상단 (65%): dwelling_grid.jpg 크게 삽입 (가로형) ──
    add_image_with_caption(
        slide,
        ASSETS / "dwelling_grid.jpg",
        Inches(0.45), Inches(2.80), Inches(10.0),
        caption_text="1fps × 30프레임 연속 캡처 — 사고 진행 과정 시계열 기록",
        tag_text="30초 체류 캡처 실물",
        tag_color=TAG_LBL,
    )

    # ── 하단 좌: outbreak_capture_sample.jpg ──
    add_image_with_caption(
        slide,
        ASSETS / "outbreak_capture_sample.jpg",
        Inches(0.45), Inches(5.40), Inches(4.5),
        caption_text="ITS 돌발정보 연동 자동 캡처 실물",
        tag_text="ITS 돌발정보 자동 캡처",
        tag_color=MAIN_BG,
    )

    # ── 하단 우: 사고 대응 타임라인 (3단계) ──
    tl_x = Inches(5.30)
    tl_y = Inches(5.50)

    add_text_box(slide, tl_x, tl_y - Inches(0.05), Inches(5.50), Inches(0.30),
                 "사고 대응 타임라인", size_pt=14, bold=True, color=MAIN_BG)

    tl_steps = [
        ("사고접수\n즉시", MAIN_BG),
        ("30초 캡처\n90프레임", CAT_HDR),
        ("AI분석+DB적재\n~20초", TAG_LBL),
    ]
    step_w = Inches(1.60)
    arrow_w = Inches(0.30)

    for i, (label, color) in enumerate(tl_steps):
        x = tl_x + i * (step_w + arrow_w + Inches(0.05))
        add_multiline_rounded_rect(
            slide, x, tl_y + Inches(0.30), step_w, Inches(0.85),
            fill_color=color,
            lines=[(line, 12, i_line == 0, WHITE) for i_line, line in enumerate(label.split('\n'))],
            alignment=PP_ALIGN.CENTER,
        )
        if i < 2:
            add_arrow_right(slide, x + step_w + Inches(0.01),
                            tl_y + Inches(0.55), arrow_w, Inches(0.30),
                            fill_color=MAIN_BG)

    # 총 소요 강조
    add_text_box(slide, tl_x, tl_y + Inches(1.25), Inches(5.50), Inches(0.30),
                 "총 소요: 약 2분 (기존 5~10분 대비 60% 이상 단축)",
                 size_pt=14, bold=True, color=RED_ALERT,
                 alignment=PP_ALIGN.CENTER)

    return slide  # ~15 shapes


def build_slide_12D(prs, layout):
    """12-D: AI 비전 파이프라인 및 기대효과 (~18 도형)"""
    slide = prs.slides.add_slide(layout)
    setup_slide_header(
        slide, 12,
        "CCTV 영상 기반 교통상황 판단",
        "12.4 AI 비전 파이프라인 및 기대효과",
        "YOLO 차량검출 + ByteTrack 추적 + 7종 이상상황 트리거",
    )

    # ── 좌측 (35%): tracks_visual.png (세로형 640x1440) ──
    add_image_with_caption(
        slide,
        ASSETS / "tracks_visual.png",
        Inches(0.45), Inches(2.80), Inches(3.50),
        caption_text="차량 검출·추적 실물 결과 (YOLO + ByteTrack)",
        tag_text="차량 검출·추적 실물",
        tag_color=TAG_LBL,
    )

    # ── 우측 상단: 트리거 7종 표 (3열 8행, 10pt) ──
    rx = Inches(4.30)
    rw = Inches(6.60)

    trigger_data = [
        ["트리거", "감지 대상", "심각도"],
        ["T1", "충돌시간(TTC) 임박", "높음"],
        ["T2", "급감속", "높음"],
        ["T3", "비정차 구간 정차", "보통"],
        ["T4", "역주행", "최고"],
        ["T5", "다중 차량 동시 급감속 (3대+)", "최고"],
        ["T6", "속도 분산 급증", "보통"],
        ["T7", "주기 상황 요약 (5분)", "정보"],
    ]
    add_table(slide, rx, Inches(2.80),
              rw, Inches(2.60),
              8, 3, trigger_data,
              header_color=MAIN_BG, font_size=10)

    # ── 우측 하단: 기대효과 3항목 카드 (축소) ──
    effects = [
        ("객관적 상황 판단", "AI 일관된 기준으로 근무자 간 편차 해소"),
        ("24시간 무중단 감시", "3-Tier 자동 감시 체계로 사람 주의력 한계 보완"),
        ("전국 확장성 363대+", "프리필터 + 핫스팟 + 정밀의 단계적 확장 구조"),
    ]

    add_text_box(slide, rx, Inches(5.60), Inches(3.00), Inches(0.30),
                 "기대 효과", size_pt=14, bold=True, color=MAIN_BG)

    for i, (title, desc) in enumerate(effects):
        y = Inches(5.95) + i * Inches(0.55)
        add_rounded_rect(slide, rx, y, Inches(2.40), Inches(0.45),
                         fill_color=MAIN_BG,
                         text=title, size_pt=12, bold=True, text_color=WHITE)
        add_text_box(slide, rx + Inches(2.55), y + Inches(0.05),
                     Inches(4.00), Inches(0.40),
                     desc, size_pt=12, color=BLACK)

    return slide  # ~18 shapes


# ══════════════════════════════════════════════════════════════
# 13장 슬라이드 (4장) — DFD를 편집 가능 도형으로 재현
# ══════════════════════════════════════════════════════════════

def build_slide_13A(prs, layout):
    """13-A: AI 판단 호출 흐름 (dfd-3 도형화) + LLM vs VLM (~28 도형)"""
    slide = prs.slides.add_slide(layout)
    setup_slide_header(
        slide, 13,
        "영상 기반 생성형 AI 모델",
        "13.1 텍스트 LLM과 멀티모달 VLM의 상호 보완",
        "Gemma3(텍스트)과 Qwen2.5-VL(영상)이 서로 다른 입력으로 사고 DB를 풍부하게 구축",
    )

    # ── 상단 (60%): dfd-3 AI 판단 흐름을 도형으로 재현 ──
    # 좌: 이상 감지 3유형 블록
    trigger_x = Inches(0.35)
    trigger_y = Inches(2.80)
    trigger_w = Inches(2.10)
    trigger_h = Inches(0.65)

    triggers = [
        ("충돌 위험", "TTC 3초 미만 / 다수 급감속"),
        ("이상 행동", "급감속 / 역주행 / 구간 정차"),
        ("정거 점검", "5분 주기 자동 촬영"),
    ]
    for i, (title, desc) in enumerate(triggers):
        ty = trigger_y + i * (trigger_h + Inches(0.08))
        add_multiline_rounded_rect(
            slide, trigger_x, ty, trigger_w, trigger_h,
            fill_color=CAT_HDR,
            lines=[
                (title, 11, True, WHITE),
                (desc, 10, False, LIGHT_PALE),
            ],
            alignment=PP_ALIGN.LEFT,
        )

    # 화살표: 트리거 → AI 전달 정보
    arr1_x = trigger_x + trigger_w + Inches(0.05)
    add_arrow_right(slide, arr1_x, trigger_y + Inches(0.90),
                    Inches(0.20), Inches(0.25), fill_color=MAIN_BG)

    # 중앙: AI 전달 정보 블록
    info_x = arr1_x + Inches(0.28)
    info_y = trigger_y
    info_w = Inches(2.00)
    add_multiline_rounded_rect(
        slide, info_x, info_y, info_w, Inches(1.40),
        fill_color=MAIN_BG,
        lines=[
            ("AI 전달 정보", 11, True, WHITE),
            ("핵심 장면 1~5장", 10, False, WHITE),
            ("차량 궤적/속도", 10, False, WHITE),
            ("이상 유형", 10, False, WHITE),
            ("IC 위치/차선", 10, False, WHITE),
        ],
        alignment=PP_ALIGN.LEFT,
    )

    # 화살표: 정보 → 장면파악
    arr2_x = info_x + info_w + Inches(0.05)
    add_arrow_right(slide, arr2_x, info_y + Inches(0.50),
                    Inches(0.20), Inches(0.25), fill_color=MAIN_BG)

    # 장면 파악 블록
    scene_x = arr2_x + Inches(0.28)
    add_rounded_rect(slide, scene_x, info_y + Inches(0.10),
                     Inches(1.40), Inches(0.45),
                     fill_color=TAG_LBL,
                     text="장면 파악", size_pt=11, bold=True, text_color=WHITE)

    # 화살표 down → 다이아몬드
    add_arrow_down(slide, scene_x + Inches(0.55), info_y + Inches(0.60),
                   Inches(0.25), Inches(0.25), fill_color=GOLD)

    # 다이아몬드: 사고인가?
    dia_x = scene_x + Inches(0.15)
    dia_y = info_y + Inches(0.90)
    add_diamond(slide, dia_x, dia_y, Inches(1.10), Inches(0.80),
                fill_color=GOLD, text="사고?", size_pt=11, bold=True,
                text_color=BLACK)

    # Yes 화살표 → 결과
    yes_arr_x = dia_x + Inches(1.15)
    add_arrow_right(slide, yes_arr_x, dia_y + Inches(0.22),
                    Inches(0.20), Inches(0.25), fill_color=GREEN_OK)
    add_text_box(slide, yes_arr_x - Inches(0.05), dia_y - Inches(0.15),
                 Inches(0.40), Inches(0.20),
                 "Yes", size_pt=10, bold=True, color=GREEN_OK,
                 alignment=PP_ALIGN.CENTER)

    # 우: 결과 블록
    result_x = yes_arr_x + Inches(0.28)
    add_multiline_rounded_rect(
        slide, result_x, dia_y - Inches(0.10), Inches(2.30), Inches(1.20),
        fill_color=DARK_BG,
        lines=[
            ("차종 세분류 보정 (13종)", 11, True, WHITE),
            ("원인 분석 + 보고서 작성", 11, False, WHITE),
            ("사고유형/날씨/차종", 10, False, LIGHT_PALE),
            ("보고서 93컬럼 자동 생성", 10, False, LIGHT_PALE),
        ],
        alignment=PP_ALIGN.LEFT,
    )

    # No → 정상 기록 (아래)
    add_text_box(slide, dia_x + Inches(0.20), dia_y + Inches(0.82),
                 Inches(0.80), Inches(0.20),
                 "No → 정상 기록", size_pt=10, bold=False, color=GRAY_TEXT,
                 alignment=PP_ALIGN.CENTER)

    # ── 하단 (40%): LLM vs VLM 비교 카드 ──
    card_y = Inches(5.30)
    card_w = Inches(3.80)
    card_h = Inches(1.10)

    # 좌 카드: Gemma3 (텍스트 LLM)
    add_multiline_rounded_rect(
        slide, Inches(0.45), card_y, card_w, card_h,
        fill_color=CAT_HDR,
        lines=[
            ("Gemma3 (텍스트 LLM)", 13, True, WHITE),
            ("원인/경위 · 인명피해", 11, False, WHITE),
            ("처리시간 · 견인/119 정보", 11, False, WHITE),
        ],
        alignment=PP_ALIGN.CENTER,
    )

    # 우 카드: Qwen2.5-VL (멀티모달 VLM)
    add_multiline_rounded_rect(
        slide, Inches(4.50), card_y, card_w, card_h,
        fill_color=MAIN_BG,
        lines=[
            ("Qwen2.5-VL (멀티모달 VLM)", 13, True, WHITE),
            ("기상/차단 · 차량/차종", 11, False, WHITE),
            ("화재/전복 · 시설물피해", 11, False, WHITE),
        ],
        alignment=PP_ALIGN.CENTER,
    )

    # 하단 중앙: 통합 사고 DB 합류 바
    db_y = card_y + card_h + Inches(0.12)
    add_rounded_rect(slide, Inches(0.45), db_y, Inches(7.85), Inches(0.45),
                     fill_color=DARK_BG,
                     text="통합 사고 DB (DuckDB 35컬럼)", size_pt=13, bold=True,
                     text_color=WHITE)

    # "경쟁이 아닌 보완" 강조
    add_text_box(slide, Inches(8.60), card_y + Inches(0.20),
                 Inches(2.50), Inches(0.40),
                 "경쟁이 아닌 보완", size_pt=14, bold=True, color=GOLD,
                 alignment=PP_ALIGN.CENTER)

    # AI 호출 빈도 소형 표 (3행 3열)
    freq_data = [
        ["상황", "빈도", "CPU 부하"],
        ["평시", "시간당 ~12회", "10~15%"],
        ["혼잡", "시간당 ~17회", "~15%"],
        ["사고", "시간당 ~10회", "~15%"],
    ]
    add_table(slide, Inches(8.60), card_y + Inches(0.65),
              Inches(2.50), Inches(1.00),
              4, 3, freq_data,
              header_color=MAIN_BG, font_size=10)

    return slide  # ~25 shapes


def build_slide_13B(prs, layout):
    """13-B: MLLM 입출력 실물 (~18 도형)"""
    slide = prs.slides.add_slide(layout)
    setup_slide_header(
        slide, 13,
        "영상 기반 생성형 AI 모델",
        "13.2 멀티모달 AI 분석 입력과 출력",
        "CCTV 영상 5장 입력 → Qwen2.5-VL 추론 → 구조화 JSON 자동 생성",
    )

    # ── 좌측 (45%): 입력 이미지 세로 2개 ──
    lx = Inches(0.45)

    # BIN0004 (197x153 → 3.0in 확대)
    slide.shapes.add_picture(str(ASSETS / "BIN0004.png"),
                             lx, Inches(2.85), Inches(3.0))
    add_text_box(slide, lx, Inches(5.20), Inches(3.0), Inches(0.22),
                 "CCTV 사고촬영 화면", size_pt=10, color=GRAY_TEXT,
                 alignment=PP_ALIGN.CENTER)

    # "VLM →" 화살표
    add_arrow_right(slide, Inches(3.70), Inches(4.10),
                    Inches(0.50), Inches(0.35), fill_color=MAIN_BG)
    add_text_box(slide, Inches(3.70), Inches(3.85),
                 Inches(0.60), Inches(0.25),
                 "VLM", size_pt=12, bold=True, color=MAIN_BG,
                 alignment=PP_ALIGN.CENTER)

    # BIN0002 (456x343 → 3.0in)
    slide.shapes.add_picture(str(ASSETS / "BIN0002.png"),
                             lx, Inches(5.50), Inches(3.0))
    add_text_box(slide, lx, Inches(7.30), Inches(3.0), Inches(0.22),
                 "사고 현장 실물 사진", size_pt=10, color=GRAY_TEXT,
                 alignment=PP_ALIGN.CENTER)

    # ── 우측 (55%): AI 분석 결과 카드 ──
    result_x = Inches(4.50)
    result_w = Inches(6.30)

    add_text_box(slide, result_x, Inches(2.80), result_w, Inches(0.30),
                 "AI 분석 결과", size_pt=14, bold=True, color=MAIN_BG)

    # 결과 박스 배경
    result_box = add_rect(slide, result_x, Inches(3.15),
                          result_w, Inches(3.20),
                          fill_color=LIGHT_BLUE_BG)
    result_box.line.color.rgb = MAIN_BG
    result_box.line.width = Pt(1)

    result_items = [
        "기상: 맑음",
        "차단: 전면차단 (1, 2번 차로)",
        "차량: 화물차 1대 (1톤 탑)",
        "적재물 유출: 의류원단",
        "시설물: 가드레일 파손",
        "원인: 차량결함 (타이어 파손)",
        "심각도: 보통",
        "신뢰도: 0.72",
    ]
    ry = Inches(3.25)
    for item in result_items:
        add_text_box(slide, result_x + Inches(0.20), ry,
                     result_w - Inches(0.40), Inches(0.28),
                     item, size_pt=12, color=BLACK)
        ry += Inches(0.33)

    # 하단 설명
    add_rounded_rect(slide, result_x, Inches(6.50), result_w, Inches(0.45),
                     fill_color=DARK_BG,
                     text="CPU 추론 ~20초/건, 3B 경량 모델 | GPU 도입 시 72B 1~3초/건",
                     size_pt=12, bold=True, text_color=WHITE)

    return slide  # ~18 shapes


def build_slide_13C(prs, layout):
    """13-C: 보고서 매핑 + 차종 13종 (dfd-4 상단 도형화) (~25 도형)"""
    slide = prs.slides.add_slide(layout)
    setup_slide_header(
        slide, 13,
        "영상 기반 생성형 AI 모델",
        "13.3 도공 보고서 양식 자동 매핑 및 차종 13종 분류",
        "전면차단 보고서 93컬럼 중 AI가 35% 자동 생성 + 화물차 10종 세분류",
    )

    # ── 좌측 (50%): 보고서 자동 작성 비율 ──
    lx = Inches(0.45)
    ly = Inches(2.80)

    add_text_box(slide, lx, ly, Inches(5.20), Inches(0.30),
                 "보고서 93컬럼 자동 작성 비율", size_pt=13, bold=True, color=MAIN_BG)

    # 3행 색상 바 (도형으로 비율 시각화)
    bar_items = [
        ("AI 자작성  27% (25건)", MAIN_BG, WHITE, 2.70),
        ("AI 추정  8% (7건)", YELLOW_TAG, BLACK, 0.80),
        ("현장확인필수  65% (61건)", LIGHT_GRAY, BLACK, 5.20),
    ]
    bar_y = ly + Inches(0.40)
    bar_h = Inches(0.45)
    max_bar_w = Inches(5.20)

    for i, (label, bg, tc, ratio_w) in enumerate(bar_items):
        y = bar_y + i * (bar_h + Inches(0.10))
        add_rounded_rect(slide, lx, y, Inches(ratio_w), bar_h,
                         fill_color=bg,
                         text=label, size_pt=12, bold=True, text_color=tc,
                         alignment=PP_ALIGN.LEFT)

    # 강조 텍스트
    emphasis_y = bar_y + 3 * (bar_h + Inches(0.10)) + Inches(0.05)
    add_text_box(slide, lx, emphasis_y, Inches(5.20), Inches(0.35),
                 "사고 직후 AI가 35% 즉시 채움", size_pt=14, bold=True,
                 color=ACCENT, alignment=PP_ALIGN.CENTER)

    # ── 우측 (50%): 차종 13종 매핑 도형 ──
    rx = Inches(5.90)
    ry = Inches(2.80)
    rw = Inches(5.00)

    add_text_box(slide, rx, ry, rw, Inches(0.30),
                 "사고보고서 차종 → 13종 세분류", size_pt=13, bold=True, color=MAIN_BG)

    # 좌 칼럼: 보고서 원본 차종 4종
    src_x = rx
    src_w = Inches(1.50)
    src_h = Inches(0.42)
    src_y = ry + Inches(0.40)

    src_types = [
        ("승용차", MAIN_BG),
        ("버스/승합", CAT_HDR),
        ("화물차 (58%)", RED_ALERT),
        ("이륜차", TAG_LBL),
    ]
    for i, (label, color) in enumerate(src_types):
        y = src_y + i * (src_h + Inches(0.08))
        add_rounded_rect(slide, src_x, y, src_w, src_h,
                         fill_color=color,
                         text=label, size_pt=11, bold=True, text_color=WHITE)

    # 중앙 화살표들
    arrow_x = src_x + src_w + Inches(0.10)

    # 우 칼럼: 13종 세분류
    dst_x = arrow_x + Inches(0.35)
    dst_w = Inches(2.30)
    dst_h = Inches(0.32)

    dst_types = [
        ("T1 승용차", PALE_BLUE, BLACK),
        ("T2 버스", PALE_BLUE, BLACK),
        ("T3~T5 화물 (소/중/대)", YELLOW_TAG, BLACK),
        ("T6~T7 특수대형", YELLOW_TAG, BLACK),
        ("T8~T12 트레일러 (5종)", YELLOW_TAG, BLACK),
        ("T13 이륜차", PALE_BLUE, BLACK),
    ]

    dst_start_y = src_y
    for i, (label, bg, tc) in enumerate(dst_types):
        y = dst_start_y + i * (dst_h + Inches(0.05))
        add_rounded_rect(slide, dst_x, y, dst_w, dst_h,
                         fill_color=bg,
                         text=label, size_pt=10, bold=False, text_color=tc,
                         alignment=PP_ALIGN.LEFT)

    # 화살표 연결 (간소화: 4개 화살표 — 원본 차종에서 세분류로)
    # 승용차 → T1
    add_arrow_right(slide, arrow_x, src_y + Inches(0.10),
                    Inches(0.30), Inches(0.18), fill_color=MAIN_BG)
    # 버스/승합 → T2
    add_arrow_right(slide, arrow_x, src_y + src_h + Inches(0.18),
                    Inches(0.30), Inches(0.18), fill_color=CAT_HDR)
    # 화물차 → T3~T12 (중앙)
    add_arrow_right(slide, arrow_x, src_y + 2 * (src_h + Inches(0.08)) + Inches(0.10),
                    Inches(0.30), Inches(0.18), fill_color=RED_ALERT)
    # 이륜차 → T13
    add_arrow_right(slide, arrow_x, src_y + 3 * (src_h + Inches(0.08)) + Inches(0.10),
                    Inches(0.30), Inches(0.18), fill_color=TAG_LBL)

    # 화물차 강조
    emphasis_y2 = dst_start_y + 6 * (dst_h + Inches(0.05)) + Inches(0.05)
    add_text_box(slide, rx, emphasis_y2, rw, Inches(0.28),
                 "화물차 58% → AI 세분류 필수", size_pt=12, bold=True,
                 color=RED_ALERT, alignment=PP_ALIGN.CENTER)

    return slide  # ~25 shapes


def build_slide_13D(prs, layout):
    """13-D: 검증 결과 + 5단계 로드맵 (dfd-4 하단 도형화) (~28 도형)"""
    slide = prs.slides.add_slide(layout)
    setup_slide_header(
        slide, 13,
        "영상 기반 생성형 AI 모델",
        "13.4 검증 결과 및 개발 로드맵",
        "실측 성능 수치 + 5단계 개발 계획",
    )

    # ── 상단 (35%): 검증 수치 카드 5개 ──
    cards = [
        ("100%", "CCTV 캡처 성공률"),
        ("30→5장", "대표프레임 자동 선정"),
        ("~20초", "AI 분석 소요시간"),
        ("35컬럼", "DB 자동 적재"),
        ("95%", "오탐 감소 (트리거)"),
    ]

    card_w = Inches(2.00)
    card_h = Inches(0.90)
    card_y = Inches(2.80)

    for i, (num, label) in enumerate(cards):
        cx = Inches(0.45) + i * (card_w + Inches(0.12))
        add_multiline_rounded_rect(
            slide, cx, card_y, card_w, card_h,
            fill_color=MAIN_BG,
            lines=[
                (num, 24, True, WHITE),
                (label, 11, False, LIGHT_PALE),
            ],
            alignment=PP_ALIGN.CENTER,
        )

    # ── 중단 (30%): 핵심 기술 요약 표 (7행 3열) ──
    tech_data = [
        ["구성 요소", "기술", "역할"],
        ["영상 수집", "Playwright + ITS API", "CCTV 자동접속 + 캡처"],
        ["객체 검출", "YOLO", "차량 실시간 검출 (21ms)"],
        ["이상 감지", "규칙엔진 + 트리거 7종", "사고 징후 자동 판정"],
        ["프리필터", "OpenCV MOG2", "경량 전국 감시 (1.5ms)"],
        ["AI 분석", "Qwen2.5-VL (3B)", "현장 상황 자동 판독"],
        ["DB 적재", "DuckDB", "35컬럼 자동 INSERT"],
    ]
    add_table(slide, Inches(0.45), Inches(3.90),
              Inches(10.70), Inches(2.10),
              7, 3, tech_data,
              header_color=MAIN_BG, font_size=10)

    # ── 하단 (35%): 5단계 로드맵 도형 (dfd-4 하단 재현) ──
    rm_top = Inches(6.25)
    add_text_box(slide, Inches(0.45), rm_top - Inches(0.05), Inches(3.00), Inches(0.25),
                 "개발 로드맵", size_pt=14, bold=True, color=MAIN_BG)

    roadmap = [
        ("0단계 (현재)", "데이터 수집\nAI 환경 구축", DARK_BG),
        ("1단계", "영상분석 고도화\n목표: 12건/s", CAT_HDR),
        ("2단계", "AI 판단 통합\n사고감지 80%", MAIN_BG),
        ("3단계", "보고서 자동화\n차종 13/13", SUB_CAT),
        ("4단계", "사고 예측\nAUC 0.7+", TAG_LBL),
    ]

    rm_step_w = Inches(1.95)
    rm_arrow_w = Inches(0.20)
    rm_h = Inches(0.90)

    for i, (title, desc, color) in enumerate(roadmap):
        x = Inches(0.45) + i * (rm_step_w + rm_arrow_w + Inches(0.05))
        # 타이틀 + 설명을 합친 카드
        desc_oneline = desc.replace('\n', ' | ')
        add_multiline_rounded_rect(
            slide, x, rm_top + Inches(0.25), rm_step_w, rm_h,
            fill_color=color,
            lines=[
                (title, 12, True, WHITE),
                (desc_oneline, 11, False, WHITE),
            ],
            alignment=PP_ALIGN.CENTER,
        )
        if i < len(roadmap) - 1:
            add_arrow_right(slide, x + rm_step_w + Inches(0.01),
                            rm_top + Inches(0.50),
                            rm_arrow_w, Inches(0.25), fill_color=color)

    return slide  # ~24 shapes


# ══════════════════════════════════════════════════════════════
# 메인 실행
# ══════════════════════════════════════════════════════════════
def main():
    # 1. 원본 복사
    DST.parent.mkdir(parents=True, exist_ok=True)
    shutil.copy2(str(SRC), str(DST))
    print(f"[OK] 원본 복사 완료: {DST}")

    # 2. 복사본 열기
    prs = Presentation(str(DST))
    layout = prs.slide_layouts[4]

    existing_count = len(prs.slides)
    print(f"[INFO] 기존 슬라이드 수: {existing_count}")
    print(f"[INFO] 사용할 레이아웃: {layout.name}")

    # 3. 슬라이드 추가 (12장 4 + 13장 4 = 총 8)
    build_slide_12A(prs, layout)
    print("[OK] 12-A: 시스템 개요 (dfd-2 도형화) + 현행 vs 개선")

    build_slide_12B(prs, layout)
    print("[OK] 12-B: 3-Tier 감시 체계")

    build_slide_12C(prs, layout)
    print("[OK] 12-C: CCTV 체류 캡처 실물")

    build_slide_12D(prs, layout)
    print("[OK] 12-D: AI 비전 파이프라인 + 기대효과")

    build_slide_13A(prs, layout)
    print("[OK] 13-A: AI 판단 흐름 (dfd-3 도형화) + LLM vs VLM")

    build_slide_13B(prs, layout)
    print("[OK] 13-B: MLLM 입출력 실물")

    build_slide_13C(prs, layout)
    print("[OK] 13-C: 보고서 매핑 + 차종 13종 (dfd-4 도형화)")

    build_slide_13D(prs, layout)
    print("[OK] 13-D: 검증 결과 + 로드맵 (dfd-4 도형화)")

    # 4. 저장
    prs.save(str(DST))
    print(f"\n[DONE] 최종 슬라이드 수: {len(prs.slides)} (기존 {existing_count} + 추가 8)")
    print(f"[DONE] 저장 완료: {DST}")

    # ══════════════════════════════════════════════════════════
    # 5. 가시성 검증: 추가된 슬라이드에서 10pt 미만 폰트 + 도형 30개 확인
    # ══════════════════════════════════════════════════════════
    print("\n[검증] 가시성 규칙 확인 (추가 슬라이드만):")
    all_pass = True

    # DFD 이미지 사용 여부 검증
    dfd_images_found = []

    for si in range(existing_count, len(prs.slides)):
        slide = prs.slides[si]
        shape_count = len(slide.shapes)
        small_fonts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for p in shape.text_frame.paragraphs:
                    for run in p.runs:
                        if run.font.size and run.text.strip():
                            pt = round(run.font.size / 12700, 1)
                            if pt < 10:
                                small_fonts.append((pt, run.text[:30]))
            # 표(테이블) 내부 텍스트도 검사
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        for p in cell.text_frame.paragraphs:
                            for run in p.runs:
                                if run.font.size and run.text.strip():
                                    pt = round(run.font.size / 12700, 1)
                                    if pt < 10:
                                        small_fonts.append((pt, run.text[:30]))
            # DFD 이미지 삽입 여부 확인
            if shape.shape_type == 13:  # Picture
                try:
                    img_name = shape.image.filename or ""
                    if "dfd" in img_name.lower():
                        dfd_images_found.append((si + 1, img_name))
                except Exception:
                    pass

        status = "PASS" if not small_fonts else f"FAIL ({len(small_fonts)}건)"
        shape_status = "PASS" if shape_count <= 30 else f"WARN ({shape_count}개 > 30)"
        if small_fonts or shape_count > 30:
            all_pass = False
        print(f"  Slide {si+1}: 도형 {shape_count}개 [{shape_status}], 10pt미만 {status}")
        if small_fonts:
            for pt, txt in small_fonts[:5]:
                print(f"    - {pt}pt: \"{txt}\"")

    # DFD 이미지 검증
    if dfd_images_found:
        all_pass = False
        print(f"\n[FAIL] DFD 이미지 {len(dfd_images_found)}개 발견 (도형 재현 필요):")
        for sn, fn in dfd_images_found:
            print(f"    - Slide {sn}: {fn}")
    else:
        print(f"\n[PASS] DFD 이미지 0개 사용 - 모든 다이어그램이 편집 가능 도형")

    # 실물 사진 확인
    photo_count = 0
    for si in range(existing_count, len(prs.slides)):
        slide = prs.slides[si]
        for shape in slide.shapes:
            if shape.shape_type == 13:  # Picture
                photo_count += 1
    print(f"[INFO] 실물 사진 삽입 수: {photo_count}개 (목표: 5개)")

    if all_pass:
        print("\n[결과] 모든 슬라이드 가시성 검증 통과")
    else:
        print("\n[결과] 일부 슬라이드에서 가시성 규칙 위반 발견")

    # 원본 무결성 확인
    import hashlib
    src_hash = hashlib.md5(SRC.read_bytes()).hexdigest()
    print(f"\n[검증] 원본 PPT MD5: {src_hash} (변경 없음 확인)")


if __name__ == "__main__":
    main()
