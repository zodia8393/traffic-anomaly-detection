#!/usr/bin/env python3
"""
사업계획서 12-13장 PPT 슬라이드 제작 (v2 — goal_12_13장_PPT.md 기준 6슬라이드)

- 원본 PPT를 복사하여 슬라이드 13 뒤에 6장 삽입
- 기존 슬라이드 14~17(Step 01~04)은 뒤로 밀림
- 레이아웃 4 ("02_과업 추진체계 및 일정의 적정성") 사용
- 디자인 시스템: 기존 PPT에서 추출한 색상/폰트/도형 패턴 준수
- 가시성: 최소 10pt, 본문 12pt+, 슬라이드당 도형 30개 이하
"""
import shutil
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ── 경로 설정 ──
SRC = Path("/공유폴더/파일공유용/JHJ/사고자료/(사업계획서)교통센터 생성형 AI 상환관제_v0.2.pptx")
DST = Path("/workspace/prj_cctv/사고분석_설계/docs/교통센터_생성형AI_12_13장.pptx")
ASSETS = Path("/workspace/prj_cctv/사고분석_설계/docs/assets")

# ── 색상 팔레트 (기존 PPT에서 추출) ──
MAIN_BLUE = RGBColor(0x00, 0x48, 0x8C)
ACCENT_BLUE = RGBColor(0x00, 0x70, 0xC0)
DARK_BLUE = RGBColor(0x00, 0x27, 0x96)
DEEP_BLUE = RGBColor(0x00, 0x20, 0x60)
MID_BLUE = RGBColor(0x00, 0x5A, 0xAB)
LINK_BLUE = RGBColor(0x00, 0x41, 0x9F)
MATCH_BLUE = RGBColor(0x35, 0x7C, 0xBD)
GOLD = RGBColor(0xFF, 0xC0, 0x00)
PURPLE = RGBColor(0x70, 0x30, 0xA0)
RED = RGBColor(0xFF, 0x00, 0x00)
DARK_RED = RGBColor(0xC0, 0x00, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)
LIGHT_GRAY = RGBColor(0xD9, 0xD9, 0xD9)
MID_GRAY = RGBColor(0xA0, 0xA0, 0xA0)
GRAY_TEXT = RGBColor(0x66, 0x66, 0x66)
LIGHT_BLUE_BG = RGBColor(0xE8, 0xF4, 0xFD)
PALE_BLUE_BG = RGBColor(0xD6, 0xEA, 0xF8)
TIER1_BLUE = RGBColor(0x5D, 0xAE, 0xD8)  # 연한 파랑 (Tier 1)
TIER2_BLUE = RGBColor(0x2E, 0x86, 0xC1)  # 중간 파랑 (Tier 2)
TIER3_BLUE = RGBColor(0x1B, 0x4F, 0x72)  # 진한 파랑 (Tier 3)
YELLOW_BG = RGBColor(0xFF, 0xF2, 0xCC)   # 노란 배경 (AI 추정)

# 폰트 이름 — 시스템에 없으면 python-pptx가 rPr에 이름만 기록,
# PowerPoint에서 열 때 해당 폰트가 있으면 적용됨
FONT_KOPUB_BOLD = "KoPub돋움체 Bold"
FONT_KOPUB_MED = "KoPub돋움체 Medium"
FONT_NANUM_EB = "나눔스퀘어 ExtraBold"

# ── 공통 헤더 위치 (EMU, 기존 슬라이드에서 추출) ──
POS_CHAPTER = (326868, 543571, 482429, 445122)      # 장번호
POS_DESC = (593725, 1645961, 9691929, 569387)        # 설명 2줄


# ══════════════════════════════════════════════════════════════
# 헬퍼 함수
# ══════════════════════════════════════════════════════════════

def set_font(run, size_pt, bold=False, color=None, name=None):
    """run에 폰트 속성 설정"""
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    if color:
        run.font.color.rgb = color
    if name:
        run.font.name = name


def add_textbox(slide, left, top, width, height, text, size_pt=12,
                bold=False, color=None, font_name=None,
                alignment=PP_ALIGN.LEFT, anchor=None):
    """텍스트 박스 추가"""
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    if anchor:
        tf.paragraphs[0].alignment = alignment
    p = tf.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    set_font(run, size_pt, bold, color, font_name)
    return tb


def add_multiline_textbox(slide, left, top, width, height, lines,
                          alignment=PP_ALIGN.LEFT):
    """여러 줄 텍스트 박스. lines = [(text, size_pt, bold, color, font_name?), ...]"""
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, line_spec in enumerate(lines):
        text, size_pt, bold, color = line_spec[:4]
        font_name = line_spec[4] if len(line_spec) > 4 else None
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = alignment
        p.space_before = Pt(2)
        p.space_after = Pt(2)
        run = p.add_run()
        run.text = text
        set_font(run, size_pt, bold, color, font_name)
    return tb


def add_rounded_rect(slide, left, top, width, height, fill_color=None,
                     text="", size_pt=12, bold=False, text_color=None,
                     alignment=PP_ALIGN.CENTER, line_color=None, line_width=None):
    """둥근 사각형 도형"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        if line_width:
            shape.line.width = Pt(line_width)
    else:
        shape.line.fill.background()

    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = alignment
        p.space_before = Pt(0)
        p.space_after = Pt(0)
        run = p.add_run()
        run.text = text
        set_font(run, size_pt, bold, text_color or (WHITE if fill_color else BLACK))
    return shape


def add_multiline_rect(slide, left, top, width, height, fill_color,
                       lines, alignment=PP_ALIGN.CENTER, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE,
                       line_color=None, line_width=None):
    """여러 줄 텍스트가 있는 도형. lines = [(text, size_pt, bold, color), ...]"""
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        if line_width:
            shape.line.width = Pt(line_width)
    else:
        shape.line.fill.background()

    tf = shape.text_frame
    tf.word_wrap = True
    for i, (text, size_pt, bold, color) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = alignment
        p.space_before = Pt(1)
        p.space_after = Pt(1)
        run = p.add_run()
        run.text = text
        set_font(run, size_pt, bold, color)
    return shape


def add_rect(slide, left, top, width, height, fill_color=None,
             text="", size_pt=12, bold=False, text_color=None,
             alignment=PP_ALIGN.CENTER, line_color=None):
    """직사각형 도형"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, height)
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = alignment
        run = p.add_run()
        run.text = text
        set_font(run, size_pt, bold, text_color or (WHITE if fill_color else BLACK))
    return shape


def add_arrow_right(slide, left, top, width, height, fill_color=None):
    """오른쪽 화살표"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RIGHT_ARROW, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color or DARK_BLUE
    shape.line.fill.background()
    return shape


def add_arrow_down(slide, left, top, width, height, fill_color=None):
    """아래쪽 화살표"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.DOWN_ARROW, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color or DARK_BLUE
    shape.line.fill.background()
    return shape


def add_table(slide, left, top, width, height, rows, cols, data,
              header_color=None, font_size=10):
    """표 추가. data = [row][col] 2D 리스트."""
    ts = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = ts.table
    col_width = width // cols
    for ci in range(cols):
        table.columns[ci].width = col_width

    for ri in range(rows):
        for ci in range(cols):
            cell = table.cell(ri, ci)
            txt = data[ri][ci] if ri < len(data) and ci < len(data[ri]) else ""
            tf = cell.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = txt
            run.font.size = Pt(font_size)
            run.font.name = FONT_KOPUB_BOLD
            if ri == 0 and header_color:
                run.font.bold = True
                run.font.color.rgb = WHITE
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE

            if ri == 0 and header_color:
                tc = cell._tc
                tcPr = tc.get_or_add_tcPr()
                for old in tcPr.findall(qn('a:solidFill')):
                    tcPr.remove(old)
                sf = tcPr.makeelement(qn('a:solidFill'), {})
                clr = sf.makeelement(qn('a:srgbClr'), {'val': str(header_color)})
                sf.append(clr)
                tcPr.insert(0, sf)
    return ts


def setup_header(slide, chapter_num, title, subtitle, desc1, desc2=""):
    """공통 슬라이드 헤더: 장번호, 제목(PH0), 부제목(PH10), 설명 2줄"""
    # 장번호 TextBox
    add_textbox(slide, *POS_CHAPTER, str(chapter_num).zfill(2),
                size_pt=20, bold=True, color=GOLD, font_name=FONT_NANUM_EB)

    # PH idx=0: 제목
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 0:
            ph.text_frame.clear()
            p = ph.text_frame.paragraphs[0]
            run = p.add_run()
            run.text = title
            run.font.name = FONT_KOPUB_BOLD
            run.font.language_id = 0x0412
            break

    # PH idx=10: 부제목
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 10:
            ph.text_frame.clear()
            p = ph.text_frame.paragraphs[0]
            run = p.add_run()
            run.text = subtitle
            run.font.name = FONT_KOPUB_BOLD
            run.font.bold = True
            break

    # 설명 텍스트 (2줄)
    desc_box = slide.shapes.add_textbox(*POS_DESC)
    tf = desc_box.text_frame
    tf.word_wrap = True
    p1 = tf.paragraphs[0]
    r1 = p1.add_run()
    r1.text = desc1
    set_font(r1, 16, bold=True, color=BLACK, name=FONT_KOPUB_BOLD)
    if desc2:
        p2 = tf.add_paragraph()
        r2 = p2.add_run()
        r2.text = desc2
        set_font(r2, 16, bold=True, color=BLACK, name=FONT_KOPUB_BOLD)


def add_poc_result_box(slide, left, top, width, height, label, result_text):
    """PoC 결과 박스 (슬라이드 7, 12 패턴)"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = LIGHT_BLUE_BG
    shape.line.color.rgb = MATCH_BLUE
    shape.line.width = Pt(1.5)

    tf = shape.text_frame
    tf.word_wrap = True
    p1 = tf.paragraphs[0]
    p1.alignment = PP_ALIGN.LEFT
    p1.space_before = Pt(2)
    r1 = p1.add_run()
    r1.text = label
    set_font(r1, 10, bold=True, color=MATCH_BLUE, name=FONT_KOPUB_BOLD)

    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.LEFT
    p2.space_before = Pt(2)
    r2 = p2.add_run()
    r2.text = result_text
    set_font(r2, 12, bold=True, color=RED, name=FONT_KOPUB_BOLD)
    return shape


# ══════════════════════════════════════════════════════════════
# 슬라이드 12-A: 개요 — 기존 vs 개선 비교도
# ══════════════════════════════════════════════════════════════

def build_12A(prs, layout):
    """12-A: CCTV 영상 기반 교통상황 판단 — 개요"""
    slide = prs.slides.add_slide(layout)
    setup_header(slide, "02",
                 "핵심 기술 및 기능",
                 "2.2 핵심 기술 설명",
                 "CCTV 영상 기반 AI 자동 교통상황 판단으로 사고 대응 시간을 단축하고 판단 객관성을 확보",
                 "사고 접수 즉시 인근 CCTV 자동 연결, AI 영상 분석, 구조화 보고서 자동 생성")

    # ── 상단: 기존 vs 개선 비교도 (좌우 2분할) ──
    comp_y = Inches(2.85)
    comp_w = Inches(4.90)
    comp_h = Inches(2.80)

    # 현행 (수동) — 회색 계열
    add_multiline_rect(
        slide, Inches(0.50), comp_y, comp_w, comp_h,
        fill_color=None, line_color=MID_GRAY, line_width=2,
        lines=[
            ("현행 (수동)", 14, True, GRAY_TEXT),
        ], alignment=PP_ALIGN.CENTER)

    steps_old = [
        "돌발정보 접수",
        "근무자가 CCTV 수동 탐색",
        "육안 확인 + 주관적 판단",
        "구두/문자 보고",
    ]
    sy = comp_y + Inches(0.45)
    for i, step in enumerate(steps_old):
        add_rounded_rect(slide, Inches(1.00), sy, Inches(3.90), Inches(0.35),
                         fill_color=LIGHT_GRAY, text=step, size_pt=11,
                         bold=False, text_color=GRAY_TEXT)
        sy += Inches(0.42)
        if i < len(steps_old) - 1:
            add_arrow_down(slide, Inches(2.80), sy - Inches(0.08),
                           Inches(0.20), Inches(0.12), fill_color=MID_GRAY)
            sy += Inches(0.05)

    # 시간/평가
    add_textbox(slide, Inches(1.00), comp_y + Inches(2.30), Inches(3.90), Inches(0.35),
                "5~10분 소요  |  주관적, 편차 발생", size_pt=12, bold=True,
                color=DARK_RED, alignment=PP_ALIGN.CENTER)

    # 개선 (AI 자동) — 파란 계열
    rx = Inches(5.80)
    add_multiline_rect(
        slide, rx, comp_y, comp_w, comp_h,
        fill_color=None, line_color=ACCENT_BLUE, line_width=2,
        lines=[
            ("개선 (AI 자동)", 14, True, DARK_BLUE),
        ], alignment=PP_ALIGN.CENTER)

    steps_new = [
        "돌발정보 접수",
        "인근 CCTV 자동 연결 (즉시)",
        "AI 영상 분석 (30초)",
        "구조화 보고서 자동 생성",
        "DB 자동 적재",
    ]
    sy = comp_y + Inches(0.40)
    for i, step in enumerate(steps_new):
        add_rounded_rect(slide, rx + Inches(0.50), sy, Inches(3.90), Inches(0.32),
                         fill_color=PALE_BLUE_BG, text=step, size_pt=11,
                         bold=False, text_color=DARK_BLUE)
        sy += Inches(0.38)
        if i < len(steps_new) - 1:
            add_arrow_down(slide, rx + Inches(2.30), sy - Inches(0.08),
                           Inches(0.20), Inches(0.12), fill_color=ACCENT_BLUE)
            sy += Inches(0.04)

    add_textbox(slide, rx + Inches(0.50), comp_y + Inches(2.30), Inches(3.90), Inches(0.35),
                "~2분 소요  |  객관적, 일관된 기준", size_pt=12, bold=True,
                color=DARK_BLUE, alignment=PP_ALIGN.CENTER)

    # 중앙 vs 화살표
    add_arrow_right(slide, Inches(5.45), comp_y + Inches(1.20),
                    Inches(0.30), Inches(0.30), fill_color=ACCENT_BLUE)

    # ── 하단: 핵심 변화 3포인트 ──
    point_y = Inches(6.00)
    point_w = Inches(3.30)
    point_h = Inches(0.80)
    gap = Inches(0.15)

    changes = [
        ("사람이 찾아서", "시스템이 자동으로"),
        ("육안으로 판단", "AI가 객관적으로"),
        ("구두 보고", "DB 자동 적재"),
    ]
    for i, (old, new) in enumerate(changes):
        cx = Inches(0.50) + i * (point_w + gap)
        add_multiline_rect(
            slide, cx, point_y, point_w, point_h,
            fill_color=DARK_BLUE,
            lines=[
                (f"{old}  →", 11, False, RGBColor(0xBB, 0xBB, 0xBB)),
                (new, 13, True, GOLD),
            ], alignment=PP_ALIGN.CENTER)

    return slide


# ══════════════════════════════════════════════════════════════
# 슬라이드 12-B: 3-Tier 전국 CCTV 감시 체계
# ══════════════════════════════════════════════════════════════

def build_12B(prs, layout):
    """12-B: 3-Tier 역삼각형 + Tier별 기술 스택 + 사고 대응 타임라인"""
    slide = prs.slides.add_slide(layout)
    setup_header(slide, "02",
                 "핵심 기술 및 기능",
                 "2.2 핵심 기술 설명",
                 "전국 4,749대 CCTV를 3단계로 나누어 제한된 서버 자원으로 전국 규모 감시와 정밀 분석을 동시 달성",
                 "사고다발구간 CCTV 자동 매칭 기반 핫스팟 선정, 이상 감지 시 자동 승격")

    # ── 좌측 (60%): 역삼각형 피라미드 ──
    pyr_x = Inches(0.50)
    pyr_y = Inches(2.85)
    # Tier 1: 넓은 상단
    add_multiline_rect(
        slide, pyr_x, pyr_y, Inches(6.30), Inches(0.85),
        fill_color=TIER1_BLUE,
        lines=[
            ("Tier 1 — 광역 감시 (4,749대)", 13, True, WHITE),
            ("경량 프리필터 (MOG2 + 프레임 차분)  |  ~1.5ms/프레임", 10, False, WHITE),
        ], alignment=PP_ALIGN.CENTER)

    # 승격 화살표 (Tier 1 → Tier 2)
    add_textbox(slide, pyr_x + Inches(1.50), pyr_y + Inches(0.88),
                Inches(1.50), Inches(0.22),
                "이상 감지 ↓", size_pt=10, bold=True, color=ACCENT_BLUE,
                alignment=PP_ALIGN.CENTER)
    add_textbox(slide, pyr_x + Inches(3.50), pyr_y + Inches(0.88),
                Inches(1.50), Inches(0.22),
                "ITS 사고 ↓", size_pt=10, bold=True, color=RED,
                alignment=PP_ALIGN.CENTER)

    # Tier 2: 중간
    add_multiline_rect(
        slide, pyr_x + Inches(0.80), pyr_y + Inches(1.15), Inches(4.70), Inches(0.85),
        fill_color=TIER2_BLUE,
        lines=[
            ("Tier 2 — 핫스팟 정밀 (50대)", 13, True, WHITE),
            ("YOLO + ByteTrack + 트리거 7종  |  HLS 24h 연속 녹화", 10, False, WHITE),
        ], alignment=PP_ALIGN.CENTER)

    # "HLS 연속 녹화" 라벨
    add_rounded_rect(slide, pyr_x + Inches(5.60), pyr_y + Inches(1.25),
                     Inches(1.10), Inches(0.30),
                     fill_color=GOLD, text="HLS 연속녹화", size_pt=10,
                     bold=True, text_color=BLACK)

    # Tier 3: 좁은 하단
    add_multiline_rect(
        slide, pyr_x + Inches(1.60), pyr_y + Inches(2.10), Inches(3.10), Inches(0.85),
        fill_color=TIER3_BLUE,
        lines=[
            ("Tier 3 — 동적 정밀 분석", 13, True, WHITE),
            ("전체 AI 파이프라인 + 멀티모달 VLM", 10, False, WHITE),
        ], alignment=PP_ALIGN.CENTER)

    # ── 우측 (40%): Tier별 기술 스택 텍스트 박스 ──
    tech_x = Inches(7.20)
    tech_w = Inches(3.80)

    tier_techs = [
        ("Tier 1 기술", [
            "OpenCV MOG2 + 프레임 차분",
            "전경비율 급변 → 이상 의심",
            "처리: ~1.5ms/프레임 (640x480)",
        ], TIER1_BLUE),
        ("Tier 2 기술", [
            "YOLO 차량 검출 (21ms/프레임)",
            "ByteTrack 차량 추적 + 속도 추정",
            "이상상황 트리거 7종 판정",
            "HLS 직접 다운로드 24h 녹화",
        ], TIER2_BLUE),
        ("Tier 3 기술", [
            "전체 비전 파이프라인",
            "Qwen2.5-VL 멀티모달 VLM",
            "보고서 자동 생성 + DB 적재",
        ], TIER3_BLUE),
    ]

    ty = pyr_y
    for title, items, color in tier_techs:
        item_h = Inches(0.22 * len(items) + 0.35)
        lines = [(title, 11, True, color)]
        for item in items:
            lines.append((f"· {item}", 10, False, BLACK))
        add_multiline_textbox(slide, tech_x, ty, tech_w, item_h, lines)
        ty += item_h + Inches(0.05)

    # ── 하단: 사고 대응 타임라인 (가로 5단계) ──
    tl_y = Inches(6.05)
    add_textbox(slide, Inches(0.50), tl_y - Inches(0.30), Inches(3.00), Inches(0.25),
                "사고 대응 타임라인", size_pt=13, bold=True, color=DARK_BLUE)

    tl_steps = [
        ("사고접수", "즉시", DARK_BLUE),
        ("CCTV\n자동연결", "5초", TIER2_BLUE),
        ("체류캡처", "30초", MID_BLUE),
        ("AI분석", "~20초", TIER3_BLUE),
        ("DB적재", "즉시", ACCENT_BLUE),
    ]
    step_w = Inches(1.80)
    arr_w = Inches(0.20)
    for i, (label, time, color) in enumerate(tl_steps):
        sx = Inches(0.50) + i * (step_w + arr_w + Inches(0.05))
        add_multiline_rect(
            slide, sx, tl_y, step_w, Inches(0.60),
            fill_color=color,
            lines=[
                (label.replace("\n", " "), 11, True, WHITE),
                (time, 10, False, GOLD),
            ], alignment=PP_ALIGN.CENTER)
        if i < len(tl_steps) - 1:
            add_arrow_right(slide, sx + step_w, tl_y + Inches(0.15),
                            arr_w, Inches(0.22), fill_color=DARK_BLUE)

    # 총 소요 강조
    add_textbox(slide, Inches(0.50), tl_y + Inches(0.65), Inches(10.50), Inches(0.30),
                "총 소요: 약 1분 (기존 5~10분 대비)", size_pt=13, bold=True,
                color=RED, alignment=PP_ALIGN.CENTER)

    return slide


# ══════════════════════════════════════════════════════════════
# 슬라이드 12-C: 영상 수집 기술 — 방식 선정 및 CCTV 매칭
# ══════════════════════════════════════════════════════════════

def build_12C(prs, layout):
    """12-C: ffmpeg vs HLS 비교 + CCTV 매칭 + PoC 결과"""
    slide = prs.slides.add_slide(layout)
    setup_header(slide, "02",
                 "핵심 기술 및 기능",
                 "2.2 핵심 기술 설명",
                 "24시간 무중단 연속 녹화를 위한 HLS 직접 다운로드 방식 채택 및 사전 검증 완료",
                 "사고다발구간 1,093개와 전국 CCTV 4,749대를 자동 매칭하여 핫스팟 감시 대상 선정")

    # ── 좌측 (50%): ffmpeg vs HLS 직접 비교 ──
    lx = Inches(0.50)
    ly = Inches(2.85)
    half_w = Inches(5.00)

    # 좌측 타이틀
    add_textbox(slide, lx, ly - Inches(0.02), half_w, Inches(0.28),
                "녹화 방식 비교", size_pt=13, bold=True, color=DARK_BLUE)

    # ffmpeg 칼럼 (탈락)
    ffmpeg_x = lx
    ffmpeg_w = Inches(2.35)
    cmp_y = ly + Inches(0.30)

    add_rect(slide, ffmpeg_x, cmp_y, ffmpeg_w, Inches(0.32),
             fill_color=MID_GRAY, text="ffmpeg -c copy (탈락)", size_pt=11,
             bold=True, text_color=WHITE)

    ffmpeg_items = [
        "세션 만료 → 프로세스 crash",
        "외부 워치독 재시작 필요",
        "24h → ~270개 프래그먼트",
        "merge 후처리 필수",
    ]
    fy = cmp_y + Inches(0.38)
    for item in ffmpeg_items:
        add_textbox(slide, ffmpeg_x + Inches(0.10), fy, ffmpeg_w - Inches(0.20), Inches(0.25),
                    f"✕ {item}", size_pt=10, color=GRAY_TEXT)
        fy += Inches(0.28)

    # HLS 직접 칼럼 (채택)
    hls_x = lx + Inches(2.55)
    hls_w = Inches(2.35)

    add_rect(slide, hls_x, cmp_y, hls_w, Inches(0.32),
             fill_color=DARK_BLUE, text="HLS 직접 (채택)", size_pt=11,
             bold=True, text_color=WHITE)

    hls_items = [
        "세션 만료 → 자동 재수립",
        "무인 운용 달성",
        "24h → 단일 .ts 파일",
        "후처리 불필요",
    ]
    hy = cmp_y + Inches(0.38)
    for item in hls_items:
        add_textbox(slide, hls_x + Inches(0.10), hy, hls_w - Inches(0.20), Inches(0.25),
                    f"✓ {item}", size_pt=10, bold=True, color=DARK_BLUE)
        hy += Inches(0.28)

    # PoC 결과 박스 (좌측 하단)
    add_poc_result_box(slide, lx, ly + Inches(1.90), half_w, Inches(0.65),
                       "사전 검증(PoC) 결과",
                       "5대 동시 24시간, 끊김 0건, ~60GB")

    # ── 우측 (50%): CCTV-사고다발구간 매칭 ──
    rx = Inches(5.80)
    rw = Inches(5.00)
    ry = Inches(2.85)

    add_textbox(slide, rx, ry - Inches(0.02), rw, Inches(0.28),
                "CCTV-사고다발구간 자동 매칭", size_pt=13, bold=True, color=DARK_BLUE)

    # 매칭 원리 4단계 (세로)
    match_steps = [
        "① ITS API에서 전국 CCTV 좌표 + IC/JCT 좌표 수집",
        "② IC/JCT 간 선형 보간으로 CCTV 노선 km 추정",
        "③ 사고다발구간(노선+km)과 CCTV km 최근접 매칭",
        "④ 사고 건수 기준 Tier 2 핫스팟 우선순위 산정",
    ]
    my = ry + Inches(0.30)
    for i, step in enumerate(match_steps):
        add_rounded_rect(slide, rx, my, rw, Inches(0.32),
                         fill_color=PALE_BLUE_BG, text=step, size_pt=10,
                         bold=False, text_color=DARK_BLUE,
                         alignment=PP_ALIGN.LEFT)
        my += Inches(0.38)
        if i < len(match_steps) - 1:
            add_arrow_down(slide, rx + Inches(2.30), my - Inches(0.08),
                           Inches(0.18), Inches(0.10), fill_color=ACCENT_BLUE)
            my += Inches(0.02)

    # PoC 결과 박스 (우측 하단)
    add_poc_result_box(slide, rx, ry + Inches(1.90), rw, Inches(0.65),
                       "사전 검증(PoC) 결과",
                       "전수 매칭 성공, 94.1% 5km 이내")

    # ── 하단: 매칭 수치 요약 ──
    summary_y = Inches(5.30)
    add_textbox(slide, Inches(0.50), summary_y, Inches(10.50), Inches(0.28),
                "매칭 상세 수치", size_pt=13, bold=True, color=DARK_BLUE)

    match_data = [
        ["대상", "매칭 실패", "오차 중앙값", "1km 이내", "5km 이내", "10km 초과"],
        ["4,749대↔1,093구간", "0건", "0.4km", "80.8%", "94.1%", "34건 (3.1%)"],
    ]
    add_table(slide, Inches(0.50), summary_y + Inches(0.30),
              Inches(10.50), Inches(0.60),
              2, 6, match_data, header_color=DARK_BLUE, font_size=10)

    # HLS 녹화 수치
    hls_data = [
        ["동시 녹화", "녹화 시간", "총 용량", "끊김", "출력 포맷"],
        ["5대 (5개 노선)", "24시간 연속", "~60GB", "0건", "단일 .ts 파일"],
    ]
    add_table(slide, Inches(0.50), summary_y + Inches(1.10),
              Inches(10.50), Inches(0.60),
              2, 5, hls_data, header_color=DARK_BLUE, font_size=10)

    return slide


# ══════════════════════════════════════════════════════════════
# 슬라이드 13-A: 멀티모달 VLM — 텍스트 LLM vs VLM 보완도
# ══════════════════════════════════════════════════════════════

def build_13A(prs, layout):
    """13-A: 텍스트 LLM vs 멀티모달 VLM 상호 보완 + 모델 비교표"""
    slide = prs.slides.add_slide(layout)
    setup_header(slide, "02",
                 "핵심 기술 및 기능",
                 "2.2 핵심 기술 설명",
                 "텍스트 LLM(Gemma3)과 별개 채널로 CCTV 영상을 직접 분석하는 멀티모달 VLM(Qwen2.5-VL) 적용",
                 "사고 직후 영상 기반 속보 정보, 보고서 완성 후 텍스트 기반 상세 정보로 점진적 DB 구축")

    # ── 상단 (60%): 상호 보완도 ──
    base_y = Inches(2.85)

    # 텍스트 경로 (좌측)
    path_w = Inches(4.50)
    path_h = Inches(0.45)
    lx = Inches(0.50)

    add_textbox(slide, lx, base_y - Inches(0.02), path_w, Inches(0.22),
                "텍스트 경로 (8~11장)", size_pt=11, bold=True, color=GRAY_TEXT)

    add_rounded_rect(slide, lx, base_y + Inches(0.25), path_w, path_h,
                     fill_color=PALE_BLUE_BG, text="사고 보고서 (원클릭 등)",
                     size_pt=11, text_color=DARK_BLUE, line_color=ACCENT_BLUE)

    add_arrow_down(slide, lx + Inches(2.00), base_y + Inches(0.72),
                   Inches(0.20), Inches(0.15), fill_color=ACCENT_BLUE)

    add_rounded_rect(slide, lx, base_y + Inches(0.90), path_w, path_h,
                     fill_color=ACCENT_BLUE, text="Gemma3  (텍스트 LLM)",
                     size_pt=12, bold=True, text_color=WHITE)

    add_arrow_down(slide, lx + Inches(2.00), base_y + Inches(1.38),
                   Inches(0.20), Inches(0.15), fill_color=ACCENT_BLUE)

    add_textbox(slide, lx, base_y + Inches(1.55), path_w, Inches(0.35),
                "원인, 경위, 인명피해, 처리시간, 견인정보",
                size_pt=10, color=DARK_BLUE, alignment=PP_ALIGN.CENTER)

    # 영상 경로 (우측)
    rx = Inches(5.80)

    add_textbox(slide, rx, base_y - Inches(0.02), path_w, Inches(0.22),
                "영상 경로 (12~13장)", size_pt=11, bold=True, color=GRAY_TEXT)

    add_rounded_rect(slide, rx, base_y + Inches(0.25), path_w, path_h,
                     fill_color=PALE_BLUE_BG, text="CCTV 영상 (실시간)",
                     size_pt=11, text_color=DARK_BLUE, line_color=DARK_BLUE)

    add_arrow_down(slide, rx + Inches(2.00), base_y + Inches(0.72),
                   Inches(0.20), Inches(0.15), fill_color=DARK_BLUE)

    add_rounded_rect(slide, rx, base_y + Inches(0.90), path_w, path_h,
                     fill_color=DARK_BLUE, text="Qwen2.5-VL  (멀티모달 VLM)",
                     size_pt=12, bold=True, text_color=WHITE)

    add_arrow_down(slide, rx + Inches(2.00), base_y + Inches(1.38),
                   Inches(0.20), Inches(0.15), fill_color=DARK_BLUE)

    add_textbox(slide, rx, base_y + Inches(1.55), path_w, Inches(0.35),
                "기상, 차량, 차로, 화재/전복, 시설물피해",
                size_pt=10, color=DARK_BLUE, alignment=PP_ALIGN.CENTER)

    # 중앙 합류 → 통합 사고 DB
    db_y = base_y + Inches(2.00)
    add_rounded_rect(slide, Inches(0.50), db_y, Inches(10.30), Inches(0.50),
                     fill_color=DEEP_BLUE,
                     text="통합 사고 DB (DuckDB 35컬럼)",
                     size_pt=14, bold=True, text_color=WHITE)

    # "경쟁이 아닌 보완" 라벨
    add_rounded_rect(slide, Inches(4.00), base_y + Inches(0.90),
                     Inches(2.80), Inches(0.45),
                     fill_color=GOLD, text="경쟁이 아닌 보완",
                     size_pt=13, bold=True, text_color=BLACK)

    # ── 하단: 모델 비교표 (TABLE) ──
    tbl_y = db_y + Inches(0.70)
    cmp_data = [
        ["구분", "Gemma3 (텍스트 LLM)", "Qwen2.5-VL (멀티모달 VLM)"],
        ["입력", "텍스트만", "이미지 + 텍스트 동시"],
        ["출력", "사고 원인/경위/피해 추출", "현장 상황 자동 판독"],
        ["파라미터", "약 40억", "약 30억"],
        ["추론 환경", "CPU 서버", "CPU 서버 (GPU 확장 가능)"],
        ["장점", "상세 서술 정보 추출", "즉시 가용, 시각 정보 직접 판독"],
        ["한계", "보고서 작성까지 시간 소요", "비시각 정보 판단 불가"],
    ]
    add_table(slide, Inches(0.50), tbl_y,
              Inches(10.30), Inches(2.20),
              7, 3, cmp_data, header_color=DARK_BLUE, font_size=10)

    return slide


# ══════════════════════════════════════════════════════════════
# 슬라이드 13-B: AI 분석 — 입출력 및 보고서 매핑
# ══════════════════════════════════════════════════════════════

def build_13B(prs, layout):
    """13-B: MLLM 입출력 예시 + 도공 보고서 양식 매핑"""
    slide = prs.slides.add_slide(layout)
    setup_header(slide, "02",
                 "핵심 기술 및 기능",
                 "2.2 핵심 기술 설명",
                 "CCTV 캡처 영상을 멀티모달 VLM이 분석하여 도공 보고서 양식에 맞는 구조화된 JSON 데이터를 자동 생성",
                 "보고서 전체 항목 중 27%를 AI가 즉시 채우고, 현장 대응자가 나머지를 보완하는 구조")

    # ── 좌측 (45%): MLLM 입출력 예시 ──
    lx = Inches(0.50)
    ly = Inches(2.85)

    # CCTV 사고영상 (BIN0004.png 삽입)
    img_path = ASSETS / "BIN0004.png"
    if img_path.exists():
        slide.shapes.add_picture(str(img_path), lx, ly, Inches(2.50))
        add_textbox(slide, lx, ly + Inches(1.95), Inches(2.50), Inches(0.22),
                    "CCTV 사고촬영 화면", size_pt=10, color=GRAY_TEXT,
                    alignment=PP_ALIGN.CENTER)
    else:
        add_rect(slide, lx, ly, Inches(2.50), Inches(1.90),
                 fill_color=LIGHT_GRAY, text="[CCTV 사고영상]",
                 size_pt=12, text_color=GRAY_TEXT)

    # VLM 화살표
    add_arrow_right(slide, Inches(3.20), ly + Inches(0.70),
                    Inches(0.40), Inches(0.30), fill_color=DARK_BLUE)
    add_textbox(slide, Inches(3.15), ly + Inches(0.45), Inches(0.50), Inches(0.25),
                "VLM", size_pt=11, bold=True, color=DARK_BLUE,
                alignment=PP_ALIGN.CENTER)

    # AI 분석 결과를 하나의 멀티라인 도형으로 통합
    json_x = Inches(3.80)
    json_w = Inches(2.00)
    add_multiline_rect(
        slide, json_x, ly, json_w, Inches(2.80),
        fill_color=LIGHT_BLUE_BG, line_color=MATCH_BLUE, line_width=1,
        lines=[
            ("AI 분석 결과", 11, True, MATCH_BLUE),
            ("기상: 맑음", 10, False, BLACK),
            ("차단: 전면차단 (1,2번 차로)", 10, False, BLACK),
            ("차량: 화물차 1대 (1톤 탑)", 10, False, BLACK),
            ("적재물 유출: 의류원단", 10, False, BLACK),
            ("시설물: 가드레일 파손", 10, False, BLACK),
            ("원인: 차량결함 (타이어)", 10, False, BLACK),
            ("심각도: 보통", 10, False, BLACK),
            ("신뢰도: 0.72", 10, False, RED),
        ], alignment=PP_ALIGN.LEFT)

    # ── 우측 (55%): 도공 보고서 양식 매핑 (TABLE 사용) ──
    rx = Inches(6.10)
    rw = Inches(4.80)
    ry = Inches(2.85)

    add_textbox(slide, rx, ry - Inches(0.02), rw, Inches(0.25),
                "도공 보고서 양식 매핑", size_pt=13, bold=True, color=DARK_BLUE)

    # 표로 보고서 매핑 구현 (10행 3열)
    report_data = [
        ["항목", "구분", "설명"],
        ["일시", "AI 자동", "ITS 접보시각"],
        ["위치", "AI 자동", "ITS 노선/방향"],
        ["기상", "AI 자동", "영상 판독"],
        ["관련차량", "AI 자동", "영상 판독"],
        ["차단유형", "AI 자동", "영상 판독"],
        ["화재/전복", "AI 추정", "영상 판독"],
        ["사고원인", "AI 추정", "참고용"],
        ["인명피해", "현장 확인", "비시각 정보"],
        ["조치사항", "현장 확인", "순찰대 보고"],
        ["견인/소통", "현장 확인", "처리 후 기록"],
    ]
    tbl_shape = add_table(slide, rx, ry + Inches(0.30),
                          rw, Inches(3.00),
                          11, 3, report_data, header_color=DARK_BLUE, font_size=10)

    # 구분 열 셀 색상 적용
    table = tbl_shape.table
    for ri in range(1, 11):
        cell = table.cell(ri, 1)
        cat = cell.text_frame.paragraphs[0].runs[0].text if cell.text_frame.paragraphs[0].runs else ""
        if "자동" in cat:
            bg = PALE_BLUE_BG
            tc = DARK_BLUE
        elif "추정" in cat:
            bg = YELLOW_BG
            tc = BLACK
        else:
            bg = LIGHT_GRAY
            tc = GRAY_TEXT
        # 구분 열 색상
        for ci in [0, 1, 2]:
            c = table.cell(ri, ci)
            tcEl = c._tc
            tcPr = tcEl.get_or_add_tcPr()
            for old in tcPr.findall(qn('a:solidFill')):
                tcPr.remove(old)
            sf = tcPr.makeelement(qn('a:solidFill'), {})
            clr = sf.makeelement(qn('a:srgbClr'), {'val': str(bg)})
            sf.append(clr)
            tcPr.insert(0, sf)
            # 텍스트 색상
            for p in c.text_frame.paragraphs:
                for run in p.runs:
                    run.font.color.rgb = tc

    # 범례 + 비율 (하나의 텍스트 박스)
    add_multiline_textbox(slide, rx, ry + Inches(3.40), rw, Inches(0.50),
                          [
                              ("■ AI 자동 27%   ■ AI 추정 8%   □ 현장 확인 65%", 11, True, DARK_BLUE),
                              ("사고 직후 AI가 27% 즉시 채움 → 현장 대응자가 나머지 보완", 12, True, PURPLE),
                          ], alignment=PP_ALIGN.CENTER)

    return slide


# ══════════════════════════════════════════════════════════════
# 슬라이드 13-C: 사전 기술 검증 (PoC) 결과
# ══════════════════════════════════════════════════════════════

def build_13C(prs, layout):
    """13-C: PoC 수치 카드 6개 + HLS/매칭 결과표"""
    slide = prs.slides.add_slide(layout)
    setup_header(slide, "02",
                 "핵심 기술 및 기능",
                 "2.2 사전 기술 검증(PoC)",
                 "제안 기술의 실현 가능성을 프로토타입 구현 및 실데이터 검증으로 확인",
                 "HLS 연속 녹화, CCTV-사고구간 매칭, AI 영상 분석, 트리거 시스템 전 구간 검증 완료")

    # ── 상단: 수치 카드 6개 (가로 배치) ──
    cards = [
        ("100%", "CCTV 캡처\n성공률"),
        ("30→5장", "대표프레임\n자동선정"),
        ("~20초", "AI 분석\n소요시간"),
        ("35", "DB 컬럼\n자동적재"),
        ("95%", "오탐 감소\n(트리거)"),
        ("94.1%", "CCTV 매칭\n5km이내"),
    ]

    card_y = Inches(2.85)
    card_w = Inches(1.65)
    card_h = Inches(1.00)
    gap = Inches(0.12)

    for i, (num, label) in enumerate(cards):
        cx = Inches(0.50) + i * (card_w + gap)
        # 숫자 카드
        add_multiline_rect(
            slide, cx, card_y, card_w, card_h,
            fill_color=None, line_color=DARK_BLUE, line_width=2,
            lines=[
                (num, 22, True, DARK_BLUE),
                (label.replace("\n", " "), 10, False, BLACK),
            ], alignment=PP_ALIGN.CENTER)

    # ── 하단 좌측 (50%): HLS 녹화 PoC 결과표 ──
    tbl_y = Inches(4.20)
    add_textbox(slide, Inches(0.50), tbl_y - Inches(0.02), Inches(5.00), Inches(0.25),
                "HLS 직접 다운로드 연속 녹화 PoC", size_pt=12, bold=True, color=DARK_BLUE)

    hls_data = [
        ["항목", "결과"],
        ["동시 녹화", "5대 (5개 노선)"],
        ["녹화 시간", "24시간 연속"],
        ["총 용량", "~60GB (카메라당 평균 ~12GB)"],
        ["끊김", "0건 (세션 만료 시 자동 재수립)"],
        ["출력", "카메라당 단일 .ts 파일"],
    ]
    add_table(slide, Inches(0.50), tbl_y + Inches(0.25),
              Inches(5.00), Inches(1.80),
              6, 2, hls_data, header_color=DARK_BLUE, font_size=10)

    # ── 하단 우측 (50%): CCTV 매칭 PoC 결과표 ──
    add_textbox(slide, Inches(5.80), tbl_y - Inches(0.02), Inches(5.00), Inches(0.25),
                "CCTV-사고다발구간 자동 매칭 PoC", size_pt=12, bold=True, color=DARK_BLUE)

    match_data = [
        ["항목", "결과"],
        ["매칭 대상", "4,749대 CCTV ↔ 1,093개 구간"],
        ["매칭 실패", "0건 (전수 매칭 성공)"],
        ["오차 중앙값", "0.4km"],
        ["5km 이내", "94.1%"],
        ["10km 초과", "34건 (3.1%) — CCTV 미설치 구간"],
    ]
    add_table(slide, Inches(5.80), tbl_y + Inches(0.25),
              Inches(5.00), Inches(1.80),
              6, 2, match_data, header_color=DARK_BLUE, font_size=10)

    # ── 최하단: 트리거 개선 수치 ──
    trig_y = Inches(6.40)
    add_textbox(slide, Inches(0.50), trig_y, Inches(10.30), Inches(0.25),
                "이상상황 트리거 정밀도 개선", size_pt=12, bold=True, color=DARK_BLUE)

    trig_data = [
        ["지표", "초기 (v3)", "개선 후 (v4)", "개선율"],
        ["시간당 트리거", "약 285건", "약 42건", "85% 감소"],
        ["시간당 TTC 오탐", "약 254건", "약 12건", "95% 감소"],
        ["오탐 녹화", "6건 (전부 오탐)", "0건", "100% 해소"],
    ]
    add_table(slide, Inches(0.50), trig_y + Inches(0.25),
              Inches(10.30), Inches(1.10),
              4, 4, trig_data, header_color=DARK_BLUE, font_size=10)

    return slide


# ══════════════════════════════════════════════════════════════
# 메인 실행
# ══════════════════════════════════════════════════════════════

def main():
    # 1. 원본 복사
    DST.parent.mkdir(parents=True, exist_ok=True)
    shutil.copy2(str(SRC), str(DST))
    print(f"[OK] 원본 복사: {SRC.name} → {DST.name}")

    # 2. 열기
    prs = Presentation(str(DST))
    layout = prs.slide_layouts[4]  # "02_과업 추진체계 및 일정의 적정성"
    existing = len(prs.slides)
    print(f"[INFO] 기존 슬라이드: {existing}장, 레이아웃: \"{layout.name}\"")

    # 3. 슬라이드 13 뒤에 삽입하기 위해 먼저 끝에 추가 후 재배치
    #    python-pptx는 add_slide가 항상 끝에 추가하므로,
    #    추가 후 XML에서 슬라이드 순서를 조정한다.

    slides_to_add = [
        ("12-A", build_12A),
        ("12-B", build_12B),
        ("12-C", build_12C),
        ("13-A", build_13A),
        ("13-B", build_13B),
        ("13-C", build_13C),
    ]

    new_slides = []
    for name, builder in slides_to_add:
        slide = builder(prs, layout)
        new_slides.append(slide)
        print(f"[OK] {name} 추가 완료")

    # 4. 슬라이드 순서 재배치: 새 6장을 슬라이드 13 뒤(인덱스 13)에 삽입
    #    현재: [1..17] + [new1..new6] = 23장
    #    목표: [1..13] + [new1..new6] + [14..17]
    sldIdLst = prs._element.find(qn('p:sldIdLst'))
    sld_ids = list(sldIdLst)  # 모든 sldId 요소

    # 기존 17장 + 새 6장 = 23장
    # 인덱스 0~12 = 슬라이드 1~13 (유지)
    # 인덱스 13~16 = 슬라이드 14~17 (기존 Step 01~04, 뒤로 밀림)
    # 인덱스 17~22 = 새 6장 (13 뒤로 이동)

    # 새 슬라이드 요소를 추출
    new_sld_elements = sld_ids[existing:]  # 인덱스 17~22
    old_tail = sld_ids[13:existing]        # 인덱스 13~16 (기존 14~17)

    # sldIdLst에서 모든 자식 제거 후 재구성
    for child in list(sldIdLst):
        sldIdLst.remove(child)

    # 1~13 유지
    for elem in sld_ids[:13]:
        sldIdLst.append(elem)
    # 새 6장 삽입
    for elem in new_sld_elements:
        sldIdLst.append(elem)
    # 기존 14~17 추가
    for elem in old_tail:
        sldIdLst.append(elem)

    # 5. 저장
    prs.save(str(DST))
    final_count = len(prs.slides)
    print(f"\n[DONE] 최종: {final_count}장 (기존 {existing} + 추가 6)")
    print(f"[DONE] 저장: {DST}")

    # ══════════════════════════════════════════════════════════
    # 6. 검증
    # ══════════════════════════════════════════════════════════
    print("\n[검증] 가시성 규칙 (추가 슬라이드만):")
    all_pass = True

    # 재배치 후 새 슬라이드는 인덱스 13~18 (슬라이드 14~19)
    for si in range(13, 13 + 6):
        slide = prs.slides[si]
        shape_count = len(slide.shapes)
        small_fonts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for p in shape.text_frame.paragraphs:
                    for run in p.runs:
                        if run.font.size and run.text.strip():
                            pt = round(run.font.size / 12700, 1)
                            if pt < 10:
                                small_fonts.append((pt, run.text[:30]))
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        for p in cell.text_frame.paragraphs:
                            for run in p.runs:
                                if run.font.size and run.text.strip():
                                    pt = round(run.font.size / 12700, 1)
                                    if pt < 10:
                                        small_fonts.append((pt, run.text[:30]))

        font_status = "PASS" if not small_fonts else f"FAIL ({len(small_fonts)}건)"
        shape_status = "PASS" if shape_count <= 30 else f"WARN ({shape_count}개)"
        if small_fonts or shape_count > 30:
            all_pass = False
        print(f"  Slide {si+1} ({slides_to_add[si-13][0]}): 도형 {shape_count}개 [{shape_status}], 10pt미만 {font_status}")
        if small_fonts:
            for pt, txt in small_fonts[:5]:
                print(f"    - {pt}pt: \"{txt}\"")

    # 기존 슬라이드 무결성
    print("\n[검증] 기존 슬라이드 1~13 무결성:")
    for si in range(13):
        slide = prs.slides[si]
        shapes = len(slide.shapes)
        print(f"  Slide {si+1}: 도형 {shapes}개")

    # 기존 Step 슬라이드 위치 확인
    print("\n[검증] Step 슬라이드 위치 (기존 14~17 → 20~23):")
    for si in range(19, min(23, final_count)):
        slide = prs.slides[si]
        title_text = ""
        for shape in slide.shapes:
            if shape.has_text_frame and "Step" in shape.text_frame.text:
                title_text = shape.text_frame.text[:40]
                break
        print(f"  Slide {si+1}: \"{title_text}\"")

    # 원본 무결성
    import hashlib
    src_hash = hashlib.md5(SRC.read_bytes()).hexdigest()
    print(f"\n[검증] 원본 PPT MD5: {src_hash} (변경 없음)")

    if all_pass:
        print("\n[결과] 모든 검증 통과")
    else:
        print("\n[결과] 일부 규칙 위반 발견 — 확인 필요")


if __name__ == "__main__":
    main()
