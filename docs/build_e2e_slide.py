"""교통센터 12-13장 PPT에 E2E 한 장 슬라이드 추가.
상단 ①데이터소스 → ②입력 → ③출력 흐름 + 하단 ④결과영상(임베드).
기존 슬라이드 디자인(헤더바·폰트·색) 유지 — 슬라이드12 복제 후 본문 교체.
"""
import copy
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

SRC = "/workspace/_inbox/교통센터_생성형AI_상황관제_12-13장 (1).pptx"
OUT = "/workspace/_inbox/교통센터_생성형AI_상황관제_12-13장_E2E추가_20260601.pptx"
ASSETS = "/workspace/prj_cctv/사고분석_설계/docs/assets"
ITS_IMG = f"{ASSETS}/its_api_spec_crop.png"
POSTER = f"{ASSETS}/vlm_poc_poster.png"
VIDEO = "/workspace/prj_cctv/사고분석_설계/output/demo_vlm_poc.mp4"

NAVY = RGBColor(0x1C, 0x26, 0x40)
GOLD = RGBColor(0xFF, 0xC0, 0x00)
BLUE = RGBColor(0x00, 0x70, 0xC0)
LGOLD = RGBColor(0xFF, 0xF3, 0xD6)
GRAY = RGBColor(0x5A, 0x66, 0x75)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LBLUE = RGBColor(0xDD, 0xEB, 0xF7)
KOPUB = "KoPub돋움체 Bold"
NSQ = "나눔스퀘어"


def dup_slide(prs, idx):
    src = prs.slides[idx]
    new = prs.slides.add_slide(src.slide_layout)
    for ph in list(new.shapes):
        ph._element.getparent().remove(ph._element)
    for sh in src.shapes:
        new.shapes._spTree.append(copy.deepcopy(sh._element))
    return new


def set_text(tf, text, font=NSQ, size=10.5, bold=True, color=NAVY, align=PP_ALIGN.LEFT):
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run(); r.text = text
    r.font.name = font; r.font.size = Pt(size); r.font.bold = bold
    r.font.color.rgb = color
    return p


def add_box(slide, x, y, w, h, fill=WHITE, line=BLUE, line_w=1.25, round=True):
    shp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if round else MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = line; shp.line.width = Pt(line_w)
    shp.shadow.inherit = False
    return shp


def add_label(slide, x, y, w, text, color=BLUE, size=12):
    """단계 라벨 (①데이터 소스 등)"""
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(0.32))
    set_text(tb.text_frame, text, font=NSQ, size=size, bold=True, color=color)
    return tb


def add_multiline(slide, x, y, w, h, lines, size=9, lead_color=NAVY):
    """여러 줄 텍스트 (key=val 형태). lines: [(text, color, bold)]"""
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.05); tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.04)
    for i, (txt, col, bd) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(2)
        r = p.add_run(); r.text = txt
        r.font.name = NSQ; r.font.size = Pt(size); r.font.bold = bd
        r.font.color.rgb = col
    return tb


def add_arrow(slide, x, y, w, h, vertical=False):
    shp = slide.shapes.add_shape(
        MSO_SHAPE.DOWN_ARROW if vertical else MSO_SHAPE.RIGHT_ARROW,
        Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid(); shp.fill.fore_color.rgb = GOLD
    shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


# ── 빌드 ──
prs = Presentation(SRC)
n_before = len(prs.slides)
slide = dup_slide(prs, 12)

# 본문 제거 (헤더/푸터/번호/섹션타이틀 유지)
from pptx.enum.shapes import MSO_SHAPE_TYPE
for sh in list(slide.shapes):
    t = sh.text_frame.text if sh.has_text_frame else ""
    if sh.shape_type == MSO_SHAPE_TYPE.TABLE:
        sh._element.getparent().remove(sh._element)
    elif sh.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
        sh._element.getparent().remove(sh._element)
    elif t.startswith("▍"):
        sh._element.getparent().remove(sh._element)

# 제목 교체 (title placeholder "13.8 핵심 기술 요약")
for sh in slide.shapes:
    if sh.has_text_frame and "핵심 기술 요약" in sh.text_frame.text:
        sh.text_frame.clear()
        set_text(sh.text_frame, "13.9 데이터 수집부터 AI 분석까지 (End-to-End)",
                 font=KOPUB, size=15, bold=True, color=NAVY)
        break

# 부제 (▍)
sub = slide.shapes.add_textbox(Inches(0.6), Inches(1.48), Inches(10.5), Inches(0.32))
set_text(sub.text_frame, "▍ ITS Open API에서 좌표 범위를 입력하면 → CCTV 목록·스트림을 수신하고 → 영상을 수집해 → AI가 상황정보를 자동 추출",
         font=NSQ, size=11.5, bold=True, color=BLUE)

# ── 상단 3박스: ① → ② → ③ ──
TOP_Y = 2.25
BOX_H = 2.20
BW = 3.30
xs = [0.55, 4.15, 7.75]

# ① 데이터 소스 — ITS 페이지 캡쳐
add_label(slide, xs[0], TOP_Y - 0.33, BW, "① 데이터 소스 — 어디에서", BLUE)
b1 = add_box(slide, xs[0], TOP_Y, BW, BOX_H)
slide.shapes.add_picture(ITS_IMG, Inches(xs[0] + 0.08), Inches(TOP_Y + 0.1),
                         width=Inches(BW - 0.16))
cap1 = slide.shapes.add_textbox(Inches(xs[0]), Inches(TOP_Y + BOX_H - 0.02), Inches(BW), Inches(0.25))
set_text(cap1.text_frame, "ITS 국가교통정보센터 CCTV 화상자료 Open API",
         font=NSQ, size=8, bold=False, color=GRAY, align=PP_ALIGN.CENTER)

# ② 입력 (요청변수)
add_label(slide, xs[1], TOP_Y - 0.33, BW, "② 입력 — 요청 변수", BLUE)
add_box(slide, xs[1], TOP_Y, BW, BOX_H, fill=LBLUE)
in_lines = [
    ("apiKey = (발급키)", NAVY, True),
    ("type = ex  (고속도로)", NAVY, True),
    ("cctvType = 1  (HLS 실시간)", NAVY, True),
    ("minX~maxX = 127.00 ~ 127.20", NAVY, True),
    ("minY~maxY = 37.30 ~ 37.50", NAVY, True),
    ("getType = json", NAVY, True),
    ("", NAVY, False),
    ("→ 좌표 범위(BBox)로 CCTV 조회", BLUE, True),
]
add_multiline(slide, xs[1] + 0.12, TOP_Y + 0.18, BW - 0.24, BOX_H - 0.3, in_lines, size=10.5)

# ③ 출력 (응답값)
add_label(slide, xs[2], TOP_Y - 0.33, BW, "③ 출력 — 응답값", BLUE)
add_box(slide, xs[2], TOP_Y, BW, BOX_H, fill=LGOLD)
out_lines = [
    ("datacount = 167  (CCTV 167대)", RGBColor(0xC0,0,0), True),
    ("cctvname = [수도권1순환] 판교..", NAVY, True),
    ("cctvurl = http://cctvsec...", NAVY, True),
    ("coordx/coordy = 경도/위도", NAVY, True),
    ("cctvformat = HLS", NAVY, True),
    ("", NAVY, False),
    ("→ 엑셀 167건 저장", BLUE, True),
    ("   (CCTV명·스트림URL·좌표)", GRAY, False),
]
add_multiline(slide, xs[2] + 0.12, TOP_Y + 0.18, BW - 0.24, BOX_H - 0.3, out_lines, size=10.5)

# 가로 화살표 ①→② , ②→③
ay = TOP_Y + BOX_H / 2 - 0.18
add_arrow(slide, 3.78, ay, 0.34, 0.36)
add_arrow(slide, 7.38, ay, 0.34, 0.36)

# 세로 화살표 ↓ (중앙)
add_arrow(slide, 5.55, TOP_Y + BOX_H + 0.02, 0.55, 0.32, vertical=True)

# ── 하단 ④ 결과 영상 (임베드) ──
VID_Y = TOP_Y + BOX_H + 0.42
add_label(slide, 0.55, VID_Y - 0.02, 6.0, "④ 결과 — AI 분석 영상 (영상 클릭 시 재생)", RGBColor(0xC0,0,0), size=12)
# poster 1280x480 (2.667:1). width 5.6 -> height 2.10
VW, VH = 4.9, 1.84
VX = 0.6
VTOP = VID_Y + 0.30
mv = slide.shapes.add_movie(VIDEO, Inches(VX), Inches(VTOP),
                            Inches(VW), Inches(VH),
                            poster_frame_image=POSTER, mime_type="video/mp4")
# ▶ 재생 아이콘 오버레이
play = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(VX + VW/2 - 0.3),
                              Inches(VTOP + VH/2 - 0.3), Inches(0.6), Inches(0.6))
play.fill.solid(); play.fill.fore_color.rgb = RGBColor(0, 0, 0)
play.line.color.rgb = WHITE; play.line.width = Pt(1.5)
tri = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, Inches(VX + VW/2 - 0.1),
                             Inches(VTOP + VH/2 - 0.13), Inches(0.22), Inches(0.26))
tri.rotation = 90
tri.fill.solid(); tri.fill.fore_color.rgb = WHITE; tri.line.fill.background()

# 영상 설명 (우측)
desc_x = VX + min(VW, 4.95) + 0.35
desc_lines = [
    ("좌측 : CCTV 영상 + 차량 검출(승용차·버스·화물차)", NAVY, True),
    ("우측 : AI CCTV 상황 정보 (최신이 상단)", NAVY, True),
    ("", NAVY, False),
    ("Qwen2.5-VL (멀티모달 VLM) · CPU 추론", BLUE, True),
    ("40초 간격 키프레임 자동 판독:", NAVY, True),
    ("  · 기상 / 차로수 / 소통상태", GRAY, False),
    ("  · 주요차종 / 이상징후(사고·정차·역주행)", GRAY, False),
    ("", NAVY, False),
    ("→ 검출 결과를 근거(grounding)로 입력해", NAVY, True),
    ("   환각 억제 + 도공 보고서 양식 자동 채움", NAVY, True),
]
add_multiline(slide, desc_x, VTOP - 0.05, 11.15 - desc_x, VH + 0.3, desc_lines, size=10)

prs.save(OUT)
print(f"저장: {OUT}")
print(f"  슬라이드 {n_before} → {len(prs.slides)} (신규 1장 추가)")
