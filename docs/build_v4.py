"""DFD PPT v4 — 용어 간략화 + DFD 단순화."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from lxml import etree
from pptx.oxml.ns import qn

C = {
    'bg':       RGBColor(0xF8, 0xF9, 0xFA),
    'white':    RGBColor(0xFF, 0xFF, 0xFF),
    'navy':     RGBColor(0x1B, 0x2A, 0x4A),
    'cyan':     RGBColor(0x08, 0x91, 0xB2),
    'amber':    RGBColor(0xF5, 0x9E, 0x0B),
    'purple':   RGBColor(0x7C, 0x3A, 0xED),
    'green':    RGBColor(0x10, 0xB9, 0x81),
    'red':      RGBColor(0xEF, 0x44, 0x44),
    'text':     RGBColor(0x1F, 0x29, 0x37),
    'sub':      RGBColor(0x6B, 0x72, 0x80),
    'border':   RGBColor(0xE5, 0xE7, 0xEB),
    'line':     RGBColor(0x94, 0xA3, 0xB8),
    'cyan_l':   RGBColor(0xE0, 0xF7, 0xFF),
    'amber_l':  RGBColor(0xFE, 0xF3, 0xC7),
    'purple_l': RGBColor(0xED, 0xE9, 0xFE),
    'green_l':  RGBColor(0xD1, 0xFA, 0xE5),
    'red_l':    RGBColor(0xFE, 0xE2, 0xE2),
    'navy_l':   RGBColor(0xDB, 0xEA, 0xFE),
    'gray_l':   RGBColor(0xF1, 0xF5, 0xF9),
    'ph0': RGBColor(0x64, 0x74, 0x8B),
    'ph1': RGBColor(0x08, 0x91, 0xB2),
    'ph2': RGBColor(0x10, 0xB9, 0x81),
    'ph3': RGBColor(0x7C, 0x3A, 0xED),
    'ph4': RGBColor(0xEF, 0x44, 0x44),
}
FONT = 'Noto Sans CJK KR'
W = Inches(13.3)
H = Inches(7.5)


def new_prs():
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H
    return prs


def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def set_bg(slide, color=None):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color or C['white']


def rect(slide, x, y, w, h, fill=None, border=None, bw=Pt(1)):
    s = slide.shapes.add_shape(1, x, y, w, h)
    if fill:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    if border:
        s.line.color.rgb = border
        s.line.width = bw
    else:
        s.line.fill.background()
    return s


def rbox(slide, x, y, w, h, fill=None, border=None, bw=Pt(1.5), cr=0.08):
    s = slide.shapes.add_shape(5, x, y, w, h)
    s.adjustments[0] = cr
    if fill:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    if border:
        s.line.color.rgb = border
        s.line.width = bw
    else:
        s.line.fill.background()
    return s


def text(shape, txt, sz=10, bold=False, color=None,
         align=PP_ALIGN.CENTER, va=MSO_ANCHOR.MIDDLE):
    if color is None:
        color = C['text']
    tf = shape.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    tf.vertical_anchor = va
    p = tf.paragraphs[0]
    p.clear()
    r = p.add_run()
    r.text = txt
    r.font.name = FONT
    r.font.size = Pt(sz)
    r.font.bold = bold
    r.font.color.rgb = color
    p.alignment = align


def mtext(shape, lines, sz=10, color=None, align=PP_ALIGN.CENTER,
          va=MSO_ANCHOR.MIDDLE, bold_first=False):
    if color is None:
        color = C['text']
    tf = shape.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    tf.vertical_anchor = va
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.clear()
        r = p.add_run()
        r.text = line
        r.font.name = FONT
        r.font.size = Pt(sz)
        r.font.bold = (bold_first and i == 0)
        r.font.color.rgb = color
        p.alignment = align


def tbox(slide, txt, x, y, w, h, sz=10, bold=False,
         color=None, align=PP_ALIGN.CENTER, va=MSO_ANCHOR.MIDDLE):
    tb = slide.shapes.add_textbox(x, y, w, h)
    text(tb, txt, sz=sz, bold=bold, color=color, align=align, va=va)
    return tb


def mtbox(slide, lines, x, y, w, h, sz=10, color=None,
          align=PP_ALIGN.CENTER, va=MSO_ANCHOR.MIDDLE, bold_first=False):
    tb = slide.shapes.add_textbox(x, y, w, h)
    mtext(tb, lines, sz=sz, color=color, align=align, va=va, bold_first=bold_first)
    return tb


def arrow(slide, x1, y1, x2, y2, color=None, width=Pt(2.0)):
    if color is None:
        color = C['line']
    conn = slide.shapes.add_connector(1, x1, y1, x2, y2)
    conn.line.color.rgb = color
    conn.line.width = width
    ln = conn.line._ln
    if ln is not None:
        tail = etree.SubElement(ln, qn('a:tailEnd'))
        tail.set('type', 'triangle')
        tail.set('w', 'med')
        tail.set('len', 'med')
    return conn


def hdr(slide, color, title):
    rect(slide, 0, 0, W, Inches(0.50), fill=color)
    tbox(slide, title, Inches(0.30), Inches(0.06), Inches(10), Inches(0.38),
         sz=15, bold=True, color=C['white'], align=PP_ALIGN.LEFT)


def panel(slide, x, y, w, h, title, fill, border, hh=Inches(0.34)):
    rbox(slide, x, y, w, h, fill=fill, border=border, bw=Pt(2))
    rect(slide, x, y, w, hh, fill=border)
    tbox(slide, title, x, y, w, hh,
         sz=10, bold=True, color=C['white'])


def pill(slide, txt, x, y, w, h, fill):
    s = rbox(slide, x, y, w, h, fill=fill, cr=0.40)
    s.line.fill.background()
    text(s, txt, sz=8, bold=True, color=C['white'])


# ═══════════════════════════════════════════
# S1: 타이틀
# ═══════════════════════════════════════════

def build_s1(prs):
    s = blank(prs)
    set_bg(s)

    stripe = [C['cyan'], C['navy'], C['amber'], C['purple']]
    sw = W // 4
    for i, c in enumerate(stripe):
        rect(s, i * sw, 0, sw, Inches(0.08), fill=c)

    CX = Inches(1.80)
    CW = Inches(9.70)

    tbox(s, 'Data Flow Diagram',
         CX, Inches(1.80), CW, Inches(0.35),
         sz=12, color=C['sub'])

    tb = s.shapes.add_textbox(CX, Inches(2.20), CW, Inches(1.00))
    text(tb, 'CCTV 교통사고 자동 분석 시스템', sz=34, bold=True, color=C['navy'])

    rect(s, (W - Inches(5)) // 2, Inches(3.30), Inches(5), Inches(0.03), fill=C['cyan'])

    tbox(s, 'CCTV 영상 AI 분석으로 사고 자동 감지 · 차종 분류 · 보고서 생성',
         CX, Inches(3.45), CW, Inches(0.40),
         sz=13, color=C['sub'])

    pills = [
        ('영상 AI', C['cyan']),
        ('차종분류 13종', C['navy']),
        ('보고서 자동화', C['amber']),
        ('사고 예측', C['purple']),
    ]
    pw = Inches(1.60)
    ph = Inches(0.32)
    gap = Inches(0.18)
    total = pw * len(pills) + gap * (len(pills) - 1)
    px = (W - total) // 2
    py = Inches(4.10)
    for lbl, col in pills:
        pill(s, lbl, px, py, pw, ph, col)
        px += pw + gap

    tbox(s, '2025-05-13   |   1단계: CPU 로컬 AI   →   2단계: GPU 서버',
         CX, Inches(4.65), CW, Inches(0.30),
         sz=9, color=C['sub'])

    metrics = [
        ('13종', '차종 분류', C['cyan']),
        ('93컬럼', '보고서 자동화', C['navy']),
        ('5단계', '개발 로드맵', C['amber']),
        ('AUC 0.7+', '사고예측 목표', C['purple']),
    ]
    cw = Inches(2.20)
    ch = Inches(0.95)
    cgap = Inches(0.28)
    tc = cw * 4 + cgap * 3
    cx0 = (W - tc) // 2
    cy = Inches(5.30)
    for num, desc, col in metrics:
        rbox(s, cx0, cy, cw, ch, fill=C['white'], border=col, bw=Pt(2.5), cr=0.12)
        tbox(s, num, cx0, cy + Inches(0.06), cw, Inches(0.50),
             sz=20, bold=True, color=col)
        tbox(s, desc, cx0, cy + Inches(0.56), cw, Inches(0.28),
             sz=9, color=C['sub'])
        cx0 += cw + cgap


# ═══════════════════════════════════════════
# S2: Level 0 DFD — 단순화
# ═══════════════════════════════════════════

def build_s2(prs):
    s = blank(prs)
    set_bg(s)
    hdr(s, C['navy'], '전체 시스템 흐름도')

    Y0 = Inches(0.62)
    M = Inches(0.16)

    # 5개 컬럼: 입력 → 영상분석 → AI판단 → 데이터저장 → 출력
    IN_W = Inches(1.60)
    VIS_W = Inches(2.30)
    ML_W = Inches(3.00)
    DB_W = Inches(2.00)
    OUT_W = Inches(1.80)
    G = Inches(0.20)

    IN_X = M
    VIS_X = IN_X + IN_W + G
    ML_X = VIS_X + VIS_W + G
    DB_X = ML_X + ML_W + G
    OUT_X = DB_X + DB_W + G

    IH = Inches(0.46)
    IG = Inches(0.08)
    PAD = Inches(0.08)

    def items(slide, x, w, y0, labels, border_col, sz=8.5):
        y = y0
        for lbl in labels:
            b = rbox(slide, x + PAD, y, w - PAD * 2, IH,
                     fill=C['white'], border=border_col, bw=Pt(0.8), cr=0.10)
            text(b, lbl, sz=sz, color=C['text'])
            y += IH + IG
        return y

    # 입력
    IN_H = Inches(2.40)
    panel(s, IN_X, Y0, IN_W, IN_H, '입력', C['green_l'], C['green'])
    items(s, IN_X, IN_W, Y0 + Inches(0.42),
          ['CCTV 영상\n(다수 IC)', '사고 보고서\n(93컬럼)', '차축 데이터\n(WIM/VDS)'],
          C['green'], sz=8)

    # 영상 분석 (Layer 1)
    VIS_H = Inches(4.70)
    panel(s, VIS_X, Y0, VIS_W, VIS_H, '① 영상 분석', C['cyan_l'], C['cyan'])
    items(s, VIS_X, VIS_W, Y0 + Inches(0.42),
          ['차량 탐지 (YOLO)', '차량 추적 (ByteTrack)',
           '차종 분류 (13종)', '속도 · 충돌위험 계산',
           '이상 상황 감지', '핵심 장면 선정'],
          C['cyan'], sz=9)

    # AI 판단 (Layer 2 = MLLM)
    ML_H = Inches(5.80)
    panel(s, ML_X, Y0, ML_W, ML_H,
          '② AI 판단 (핵심)', C['navy_l'], C['navy'])
    ml_items = [
        ('장면 파악', False),
        ('사고 여부 판정', False),
        ('차종 세분류 보정 (13종)', True),
        ('사고 원인 추론', False),
        ('보고서 자동 작성 (93컬럼)', True),
        ('위험도 평가', False),
    ]
    my = Y0 + Inches(0.42)
    for lbl, key in ml_items:
        b = rbox(s, ML_X + PAD, my, ML_W - PAD * 2, IH,
                 fill=C['white'], border=C['navy'],
                 bw=Pt(2.0 if key else 0.8), cr=0.10)
        text(b, lbl, sz=9.5, bold=key,
             color=C['navy'] if key else C['text'])
        my += IH + IG

    poc = rbox(s, ML_X + PAD, my + Inches(0.10),
               ML_W - PAD * 2, Inches(0.50),
               fill=C['navy'], cr=0.12)
    poc.line.fill.background()
    mtext(poc, ['1단계: Qwen2.5-VL-7B (CPU 로컬)',
                '2단계: 72B 모델 (GPU 서버)'],
          sz=8, color=C['white'])

    # 데이터 저장 (Layer 3)
    DB_H = Inches(3.00)
    panel(s, DB_X, Y0, DB_W, DB_H, '③ 데이터 저장', C['amber_l'], C['amber'])
    items(s, DB_X, DB_W, Y0 + Inches(0.42),
          ['차량 궤적', '사고 기록', 'AI 분석 결과', '교통량 통계'],
          C['amber'], sz=8.5)

    # 사고 예측 (Layer 4) — 데이터저장 아래에 배치
    PRED_Y = Y0 + DB_H + Inches(0.20)
    PRED_H = Inches(2.56)
    panel(s, DB_X, PRED_Y, DB_W, PRED_H, '④ 사고 예측', C['purple_l'], C['purple'])
    items(s, DB_X, DB_W, PRED_Y + Inches(0.42),
          ['과거 데이터 학습', 'AI 위험도 활용', '위험 구간 예측', '목표: AUC 0.7+'],
          C['purple'], sz=8.5)

    # 출력
    OUT_H = Inches(5.80)
    panel(s, OUT_X, Y0, OUT_W, OUT_H, '출력', C['red_l'], C['red'])
    items(s, OUT_X, OUT_W, Y0 + Inches(0.42),
          ['사고 알림', '차종별 교통량', '원인 분석\n리포트',
           '자동 보고서\n(93컬럼)', '위험도 예측', '위험 구간\n히트맵'],
          C['red'], sz=8)

    # 화살표
    mid = Y0 + Inches(2.00)
    arrow(s, IN_X + IN_W, mid, VIS_X, mid, color=C['green'], width=Pt(2.5))
    tbox(s, '영상', IN_X + IN_W - Inches(0.05), mid - Inches(0.22),
         G + Inches(0.1), Inches(0.20), sz=7, color=C['green'])

    arrow(s, VIS_X + VIS_W, mid, ML_X, mid, color=C['cyan'], width=Pt(2.8))
    tbox(s, '이상 감지 시\n핵심 장면 전달', VIS_X + VIS_W - Inches(0.1), mid - Inches(0.30),
         G + Inches(0.4), Inches(0.28), sz=7, color=C['cyan'])

    # AI → 저장
    arrow(s, ML_X + ML_W, Y0 + Inches(1.50),
          DB_X, Y0 + Inches(1.50), color=C['navy'], width=Pt(2.0))
    tbox(s, '저장', ML_X + ML_W, Y0 + Inches(1.28),
         Inches(0.50), Inches(0.20), sz=7, color=C['navy'])

    # 저장 → AI (검색)
    arrow(s, DB_X, Y0 + Inches(2.10),
          ML_X + ML_W, Y0 + Inches(2.10), color=C['amber'], width=Pt(1.8))
    tbox(s, '유사 사례', ML_X + ML_W + Inches(0.02), Y0 + Inches(2.12),
         Inches(0.80), Inches(0.20), sz=7, color=C['amber'], align=PP_ALIGN.LEFT)

    # AI → 출력
    arrow(s, ML_X + ML_W, mid + Inches(0.60),
          OUT_X, mid + Inches(0.60), color=C['navy'], width=Pt(2.5))

    # AI → 예측
    arrow(s, ML_X + ML_W, Y0 + Inches(3.80),
          DB_X, PRED_Y + Inches(0.60), color=C['purple'], width=Pt(1.8))
    tbox(s, '위험도', ML_X + ML_W, Y0 + Inches(3.58),
         Inches(0.60), Inches(0.20), sz=7, color=C['purple'])

    # 예측 → 출력
    arrow(s, DB_X + DB_W, PRED_Y + Inches(1.20),
          OUT_X, Y0 + Inches(4.40), color=C['purple'], width=Pt(1.8))

    # 범례
    ly = Inches(6.90)
    legend = [
        (C['cyan'], '영상 분석'), (C['navy'], 'AI 판단'),
        (C['amber'], '데이터 저장'), (C['purple'], '사고 예측'),
        (C['green'], '입력'), (C['red'], '출력'),
    ]
    lx = Inches(0.40)
    for col, name in legend:
        dot = s.shapes.add_shape(9, lx, ly + Inches(0.05), Inches(0.14), Inches(0.14))
        dot.fill.solid()
        dot.fill.fore_color.rgb = col
        dot.line.fill.background()
        tbox(s, name, lx + Inches(0.20), ly, Inches(1.20), Inches(0.26),
             sz=8, color=C['sub'], align=PP_ALIGN.LEFT)
        lx += Inches(1.60)


# ═══════════════════════════════════════════
# S3: Level 1 DFD — 트리거 3그룹 + 태스크
# ═══════════════════════════════════════════

def build_s3(prs):
    s = blank(prs)
    set_bg(s)
    hdr(s, C['cyan'], 'AI 판단 호출 흐름 — 이상 감지 시에만 호출')

    Y0 = Inches(0.60)
    PAD = Inches(0.08)

    # A. 이상 감지 조건 (3그룹으로 압축)
    TRG_X = Inches(0.16)
    TRG_W = Inches(2.60)
    TRG_H = Inches(4.60)
    panel(s, TRG_X, Y0, TRG_W, TRG_H,
          '이상 감지 조건 (3유형)', C['cyan_l'], C['cyan'])

    groups = [
        ('충돌 위험', ['충돌 예상 시간 3초 미만', '다수 차량 동시 급감속'],
         C['red'], Inches(1.10)),
        ('이상 행동', ['급감속 · 급정거', '비정차 구간 정차', '역주행'],
         C['amber'], Inches(1.30)),
        ('정기 점검', ['5분 주기 자동 촬영'],
         C['cyan'], Inches(0.80)),
    ]
    gy = Y0 + Inches(0.42)
    for gname, gitems, gcol, gh in groups:
        gb = rbox(s, TRG_X + Inches(0.08), gy,
                  TRG_W - Inches(0.16), gh,
                  fill=C['white'], border=gcol, bw=Pt(1.5), cr=0.10)
        lines = [gname] + [f'  · {it}' for it in gitems]
        mtbox(s, lines, TRG_X + Inches(0.14), gy + Inches(0.04),
              TRG_W - Inches(0.28), gh - Inches(0.08),
              sz=8.5, color=C['text'], align=PP_ALIGN.LEFT,
              va=MSO_ANCHOR.TOP, bold_first=True)
        gy += gh + Inches(0.12)

    # B. AI 입력 정보
    PKG_X = TRG_X + TRG_W + Inches(0.30)
    PKG_W = Inches(2.40)
    PKG_H = Inches(3.00)
    panel(s, PKG_X, Y0, PKG_W, PKG_H,
          'AI에 전달하는 정보', C['navy_l'], C['navy'])
    pkg = ['핵심 장면 1~5장', '차량 궤적 · 속도', '이상 유형',
           'IC 위치 · 차선', '유사 과거 사고']
    py = Y0 + Inches(0.42)
    for item in pkg:
        b = rbox(s, PKG_X + PAD, py, PKG_W - PAD * 2, Inches(0.42),
                 fill=C['white'], border=C['navy'], bw=Pt(0.8), cr=0.09)
        text(b, item, sz=8.5, color=C['navy'])
        py += Inches(0.50)

    # 트리거 → 패키지 화살표
    arrow(s, TRG_X + TRG_W, Y0 + TRG_H / 2,
          PKG_X, Y0 + PKG_H / 2,
          color=C['cyan'], width=Pt(2.5))

    # C. 판단 흐름
    FX = PKG_X + PKG_W + Inches(0.35)
    TW = Inches(2.60)
    TH = Inches(0.58)
    TGAP = Inches(0.24)

    def task(slide, y, title, border=C['navy']):
        b = rbox(slide, FX, y, TW, TH,
                 fill=C['white'], border=border, bw=Pt(2.0), cr=0.10)
        text(b, title, sz=10, bold=True, color=border)
        return b

    # 패키지 → 판단1
    T1Y = Y0 + Inches(0.10)
    arrow(s, PKG_X + PKG_W, Y0 + Inches(1.00),
          FX, T1Y + TH / 2, color=C['navy'], width=Pt(2.2))

    task(s, T1Y, '장면 파악')
    arrow(s, FX + TW / 2, T1Y + TH, FX + TW / 2, T1Y + TH + TGAP,
          color=C['navy'], width=Pt(1.8))

    # 다이아몬드
    DY = T1Y + TH + TGAP
    DW = Inches(1.80)
    DH = Inches(0.65)
    DX = FX + (TW - DW) / 2
    dia = s.shapes.add_shape(4, DX, DY, DW, DH)
    dia.fill.solid()
    dia.fill.fore_color.rgb = C['amber_l']
    dia.line.color.rgb = C['amber']
    dia.line.width = Pt(2.5)
    text(dia, '사고인가?', sz=11, bold=True, color=C['amber'])

    # NO
    no_x = DX + DW + Inches(0.15)
    no_y = DY + (DH - Inches(0.42)) / 2
    arrow(s, DX + DW, DY + DH / 2, no_x, no_y + Inches(0.21),
          color=C['sub'], width=Pt(1.5))
    tbox(s, 'NO', DX + DW + Inches(0.02), DY + DH / 2 - Inches(0.22),
         Inches(0.30), Inches(0.20), sz=9, bold=True, color=C['sub'])
    nb = rbox(s, no_x, no_y, Inches(1.40), Inches(0.42),
              fill=C['gray_l'], border=C['sub'], bw=Pt(1), cr=0.12)
    text(nb, '정상 기록', sz=8, color=C['sub'])

    # YES
    T3Y = DY + DH + TGAP
    arrow(s, FX + TW / 2, DY + DH, FX + TW / 2, T3Y,
          color=C['green'], width=Pt(2.0))
    tbox(s, 'YES', FX + TW / 2 + Inches(0.06), DY + DH,
         Inches(0.40), Inches(0.22), sz=9, bold=True, color=C['green'],
         align=PP_ALIGN.LEFT)

    task(s, T3Y, '차종 세분류 보정 (13종)', border=C['cyan'])

    T4Y = T3Y + TH + TGAP
    arrow(s, FX + TW / 2, T3Y + TH, FX + TW / 2, T4Y,
          color=C['navy'], width=Pt(1.8))
    task(s, T4Y, '원인 분석 + 보고서 작성', border=C['amber'])

    # 결과 박스
    RY = T4Y + TH + TGAP
    rb = rbox(s, FX, RY, TW, Inches(0.90),
              fill=C['navy'], cr=0.10)
    rb.line.fill.background()
    arrow(s, FX + TW / 2, T4Y + TH, FX + TW / 2, RY,
          color=C['navy'], width=Pt(1.8))
    mtext(rb, ['분석 결과 (자동 생성)',
               '사고 유형 · 심각도 · 차종',
               '원인 추론 · 보고서 93컬럼'],
          sz=8.5, color=C['white'], bold_first=True)

    # D. CPU 부하 — 우측 간단 카드
    CPU_X = Inches(9.50)
    CPU_Y = Inches(3.80)
    CPU_W = Inches(3.55)
    CPU_H = Inches(2.40)
    panel(s, CPU_X, CPU_Y, CPU_W, CPU_H,
          'AI 호출 빈도', C['gray_l'], C['sub'])

    rows = [
        ('평시', '시간당 ~12회', '10~15%'),
        ('혼잡', '시간당 ~17회', '~15%'),
        ('사고', '시간당 ~10회', '~15%'),
    ]
    cxs = [CPU_X + Inches(0.10), CPU_X + Inches(0.90), CPU_X + Inches(2.10)]
    cws = [Inches(0.76), Inches(1.16), Inches(1.15)]
    rh = Inches(0.44)
    ry = CPU_Y + Inches(0.40)

    for cx, cw, hd in zip(cxs, cws, ['상황', '빈도', 'CPU 부하']):
        hb = rect(s, cx, ry, cw, rh - Inches(0.04), fill=C['sub'])
        tbox(s, hd, cx, ry, cw, rh - Inches(0.04),
             sz=8, bold=True, color=C['white'])
    ry += rh
    for ri, row in enumerate(rows):
        fill = C['white'] if ri % 2 == 0 else C['gray_l']
        for cx, cw, cell in zip(cxs, cws, row):
            rect(s, cx, ry, cw, rh - Inches(0.04), fill=fill,
                 border=C['border'], bw=Pt(0.5))
            tbox(s, cell, cx, ry, cw, rh - Inches(0.04),
                 sz=8, color=C['text'])
        ry += rh

    tbox(s, '※ 매 프레임 호출이 아닌 이상 감지 시에만 AI 호출',
         CPU_X, CPU_Y + CPU_H - Inches(0.30), CPU_W, Inches(0.26),
         sz=7.5, bold=True, color=C['red'])

    # 우상단 설명
    tbox(s, 'AI는 이상 감지 시에만 호출\n→ CPU 부하 최소화',
         CPU_X, Y0, CPU_W, Inches(3.60),
         sz=9, color=C['sub'], align=PP_ALIGN.LEFT, va=MSO_ANCHOR.BOTTOM)


# ═══════════════════════════════════════════
# S4: 도메인 + 로드맵
# ═══════════════════════════════════════════

def build_s4(prs):
    s = blank(prs)
    set_bg(s)
    hdr(s, C['amber'], '차종 분류 · 보고서 자동화 · 개발 계획')

    Y0 = Inches(0.60)
    PAD = Inches(0.08)

    # A. 차종 매핑 (5그룹으로 압축)
    MAP_X = Inches(0.16)
    MAP_W = Inches(4.20)
    MAP_H = Inches(4.80)
    panel(s, MAP_X, Y0, MAP_W, MAP_H,
          '사고보고서 차종 → 13종 세분류', C['navy_l'], C['navy'])

    # 보고서 차종 (좌)
    tbox(s, '보고서', MAP_X + Inches(0.10), Y0 + Inches(0.40),
         Inches(1.10), Inches(0.20), sz=8, bold=True, color=C['navy'])

    report = [
        ('승용차', C['cyan'], Inches(0.50)),
        ('버스/승합', C['cyan'], Inches(0.50)),
        ('화물차\n(58%!)', C['red'], Inches(1.40)),
        ('이륜차', C['purple'], Inches(0.50)),
    ]
    ry = Y0 + Inches(0.64)
    rmids = []
    for rname, col, rh in report:
        b = rbox(s, MAP_X + Inches(0.10), ry, Inches(1.10), rh,
                 fill=C['white'], border=col, bw=Pt(1.8), cr=0.10)
        text(b, rname, sz=9, bold=True, color=col)
        rmids.append(ry + rh / 2)
        ry += rh + Inches(0.12)

    # 13종 (우) — 5그룹
    tbox(s, '13종 세분류', MAP_X + Inches(2.20), Y0 + Inches(0.40),
         Inches(1.80), Inches(0.20), sz=8, bold=True, color=C['navy'])

    groups_t = [
        ('T1 승용차', C['cyan'], Inches(0.40)),
        ('T2 버스', C['cyan'], Inches(0.40)),
        ('T3~T5 화물 (소·중·대)', C['amber'], Inches(0.50)),
        ('T6~T7 특수 대형', C['amber'], Inches(0.45)),
        ('T8~T12 트레일러 (5종)', C['red'], Inches(0.50)),
        ('T13 이륜차', C['purple'], Inches(0.40)),
    ]
    ty = Y0 + Inches(0.64)
    tmids = []
    for tname, col, th in groups_t:
        b = rbox(s, MAP_X + Inches(2.20), ty, Inches(1.80), th,
                 fill=C['white'], border=col, bw=Pt(1.0), cr=0.09)
        text(b, tname, sz=8.5, bold=True, color=col)
        tmids.append(ty + th / 2)
        ty += th + Inches(0.06)

    # 매핑 화살표
    arrow(s, MAP_X + Inches(1.20), rmids[0],
          MAP_X + Inches(2.20), tmids[0], color=C['cyan'], width=Pt(1.5))
    arrow(s, MAP_X + Inches(1.20), rmids[1],
          MAP_X + Inches(2.20), tmids[1], color=C['cyan'], width=Pt(1.5))
    # 화물 → T3~T12 (3개 그룹)
    cargo_mid = (tmids[2] + tmids[4]) / 2
    arrow(s, MAP_X + Inches(1.20), rmids[2],
          MAP_X + Inches(2.20), cargo_mid, color=C['red'], width=Pt(2.0))
    arrow(s, MAP_X + Inches(1.20), rmids[3],
          MAP_X + Inches(2.20), tmids[5], color=C['purple'], width=Pt(1.5))

    # AI 판단 순서
    tbox(s, 'AI 판단 순서: 단위수 → 차축수 → 차체형태 → 크기',
         MAP_X + Inches(0.08), Y0 + MAP_H - Inches(0.62),
         MAP_W - Inches(0.16), Inches(0.26),
         sz=8, color=C['navy'], align=PP_ALIGN.CENTER)
    tbox(s, '★ 화물차 = 10종 → AI 세분류 필수',
         MAP_X + Inches(0.08), Y0 + MAP_H - Inches(0.34),
         MAP_W - Inches(0.16), Inches(0.28),
         sz=9, bold=True, color=C['red'], align=PP_ALIGN.CENTER)

    # B. 93컬럼 자동화 (우상)
    AUTO_X = MAP_X + MAP_W + Inches(0.24)
    AUTO_Y = Y0
    AUTO_W = Inches(4.20)
    AUTO_H = Inches(2.20)
    panel(s, AUTO_X, AUTO_Y, AUTO_W, AUTO_H,
          '보고서 93컬럼 자동 작성 비율', C['amber_l'], C['amber'])

    bars = [
        ('AI 자동 작성', 0.27, '27% (25건)', C['green']),
        ('AI 추정', 0.08, '8% (7건)', C['amber']),
        ('현장 확인 필수', 0.65, '65% (61건)', C['red']),
    ]
    BLW = Inches(1.30)
    BMW = AUTO_W - BLW - Inches(1.00)
    BH = Inches(0.32)
    BX = AUTO_X + BLW + Inches(0.08)
    by = AUTO_Y + Inches(0.44)
    for lbl, ratio, val, col in bars:
        tbox(s, lbl, AUTO_X + Inches(0.06), by, BLW, BH,
             sz=8, color=C['text'], align=PP_ALIGN.RIGHT, va=MSO_ANCHOR.MIDDLE)
        rect(s, BX, by + Inches(0.04), BMW, BH - Inches(0.08), fill=C['border'])
        aw = int(BMW * ratio)
        if aw > 0:
            rect(s, BX, by + Inches(0.04), aw, BH - Inches(0.08), fill=col)
        tbox(s, val, BX + aw + Inches(0.04), by, Inches(0.86), BH,
             sz=8, bold=True, color=col, align=PP_ALIGN.LEFT, va=MSO_ANCHOR.MIDDLE)
        by += BH + Inches(0.14)

    tbox(s, '흐름: 사고 감지 → 27% 자동 → 현장 알림 → 65% 보완',
         AUTO_X + Inches(0.06), AUTO_Y + AUTO_H - Inches(0.34),
         AUTO_W - Inches(0.12), Inches(0.30),
         sz=8, color=C['sub'])

    # C. 검증 패턴 (우하)
    VX = AUTO_X
    VY = AUTO_Y + AUTO_H + Inches(0.16)
    VW = AUTO_W
    VH = Inches(2.42)
    panel(s, VX, VY, VW, VH,
          '검증된 방식 — 같은 구조, 다른 입력', C['gray_l'], C['sub'])

    half = (VW - Inches(0.52)) / 2
    compare = [
        ('기존 (텍스트)', ['PDF 원문', '규칙 추출', 'AI 보정'], C['sub']),
        ('본 시스템 (영상)', ['CCTV 영상', '영상 분석', 'AI 판단'], C['navy']),
    ]
    for ci, (title, items, col) in enumerate(compare):
        cx = VX + Inches(0.08) + ci * (half + Inches(0.36))
        tbox(s, title, cx, VY + Inches(0.38), half, Inches(0.26),
             sz=9, bold=True, color=col)
        iy = VY + Inches(0.68)
        for k, item in enumerate(items):
            ib = rbox(s, cx, iy, half, Inches(0.38),
                      fill=C['white'], border=col, bw=Pt(1.0), cr=0.12)
            text(ib, item, sz=8.5, color=C['text'])
            if k < len(items) - 1:
                arrow(s, cx + half / 2, iy + Inches(0.38),
                      cx + half / 2, iy + Inches(0.48),
                      color=col, width=Pt(1.5))
            iy += Inches(0.50)

    eq_x = VX + half + Inches(0.12)
    tbox(s, '=', eq_x, VY + VH / 2 - Inches(0.15),
         Inches(0.28), Inches(0.36), sz=24, bold=True, color=C['sub'])

    # D. 로드맵 (하단)
    RY = Inches(5.62)
    RH = Inches(1.68)
    RX = Inches(0.16)
    RW = W - Inches(0.32)
    panel(s, RX, RY, RW, RH, '개발 로드맵 — 5단계', C['gray_l'], C['border'],
          hh=Inches(0.28))

    phases = [
        ('0단계', '데이터 수집\nAI 환경 구축', '영상 50건\nAI 서빙 확인', C['ph0']),
        ('1단계', '영상 분석\n고도화', '추적 정확도\n속도 오차 15%', C['ph1']),
        ('2단계', 'AI 판단 통합\n사고 감지', '감지 정확도\n80% 이상', C['ph2']),
        ('3단계', '보고서 자동화\n데이터 축적', '자동 작성 85%\n품질 3.5/5', C['ph3']),
        ('4단계', '사고 예측\n모델 개발', 'AUC 0.7+', C['ph4']),
    ]
    N = len(phases)
    PG = Inches(0.16)
    PW = (RW - Inches(0.26) - PG * (N - 1)) / N
    PH = Inches(1.16)
    px = RX + Inches(0.13)
    py = RY + Inches(0.32)

    for i, (pname, pdesc, gate, col) in enumerate(phases):
        if i > 0:
            arrow(s, px - PG + Inches(0.02), py + PH / 2,
                  px, py + PH / 2, color=C['sub'], width=Pt(1.5))
        rbox(s, px, py, PW, PH, fill=C['white'], border=col, bw=Pt(2.0), cr=0.10)
        rect(s, px, py, PW, Inches(0.26), fill=col)
        tbox(s, pname, px, py, PW, Inches(0.26),
             sz=9, bold=True, color=C['white'])
        tbox(s, pdesc, px + Inches(0.04), py + Inches(0.28),
             PW - Inches(0.08), Inches(0.48), sz=8, color=C['text'])
        tbox(s, gate, px + Inches(0.04), py + Inches(0.78),
             PW - Inches(0.08), Inches(0.36), sz=7, color=C['sub'])
        px += PW + PG


# ═══════════════════════════════════════════

def main():
    prs = new_prs()
    build_s1(prs)
    build_s2(prs)
    build_s3(prs)
    build_s4(prs)
    out = '/workspace/prj_cctv/new/멀티모달_교통사고_DFD_v4_20250513.pptx'
    prs.save(out)
    print(f'Saved: {out}')
    return out


if __name__ == '__main__':
    main()
